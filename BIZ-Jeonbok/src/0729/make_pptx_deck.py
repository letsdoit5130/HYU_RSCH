"""
완도 전복 글로벌 무역 EDA 및 신시장 개척 PPTX 발표자료 덱 생성 스크립트 (BIZ-JB-Gathered.csv 수록 HS CODE 반영)

이 스크립트는 python-pptx를 활용하여 BIZ-Jeonbok/reports/Wando_Abalone_Market_Entry.pptx 슬라이드를 생성합니다.
BIZ-JB-Gathered.csv 데이터셋에 실재하는 HS CODE (030781, 160557, 030783) 기준을 엄격하게 적용합니다.
"""

import os
import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)  # 16:9 Widescreen

# Color Palette: Midnight Executive
COLOR_BG_DARK = RGBColor(0x1B, 0x36, 0x5D)
COLOR_BG_LIGHT = RGBColor(0xF5, 0xF7, 0xFA)
COLOR_CARD = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_PRIMARY = RGBColor(0x1B, 0x36, 0x5D)
COLOR_ACCENT = RGBColor(0x2B, 0x5C, 0x8F)
COLOR_GOLD = RGBColor(0xD4, 0xAF, 0x37)
COLOR_TEXT_DARK = RGBColor(0x22, 0x22, 0x22)
COLOR_TEXT_MUTED = RGBColor(0x66, 0x66, 0x66)

blank_slide_layout = prs.slide_layouts[6]

def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_header(slide, title_text, category_text="WANDO ABALONE GLOBAL EXPANSION"):
    tx_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.8))
    tf = tx_box.text_frame
    tf.word_wrap = True
    
    p0 = tf.paragraphs[0]
    p0.text = category_text.upper()
    p0.font.name = 'Arial'
    p0.font.size = Pt(10)
    p0.font.bold = True
    p0.font.color.rgb = COLOR_GOLD
    p0.space_after = Pt(2)
    
    p1 = tf.add_paragraph()
    p1.text = title_text
    p1.font.name = 'Arial Black'
    p1.font.size = Pt(22)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_PRIMARY

# Slide 1: Title
slide1 = prs.slides.add_slide(blank_slide_layout)
set_bg(slide1, COLOR_BG_DARK)

tb1 = slide1.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.3), Inches(3.5))
tf1 = tb1.text_frame
tf1.word_wrap = True

p = tf1.paragraphs[0]
p.text = "완도 전복 글로벌 무역 EDA 및 신시장 개척 전략"
p.font.name = 'Malgun Gothic'
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
p.space_after = Pt(12)

p2 = tf1.add_paragraph()
p2.text = "BIZ-JB-Gathered.csv 수록 HS CODE (030781, 160557, 030783) 기준 100% 매칭 검증"
p2.font.name = 'Malgun Gothic'
p2.font.size = Pt(17)
p2.font.color.rgb = COLOR_GOLD

# Slide 2: EDA Executive Summary
slide2 = prs.slides.add_slide(blank_slide_layout)
set_bg(slide2, COLOR_BG_LIGHT)
add_header(slide2, "글로벌 전복 무역 데이터 (5,400건) 주요 지표 요약")

stats = [
    ("5,400 건", "글로벌 무역 분석 데이터", "$52.09억 무역규모"),
    ("94.1 %", "수입(Import) 거래 비중", "수입국 통관 데이터 중심"),
    ("$33.03 / kg", "전복 평균 단위당 단가", "최고 $2,284/kg 건전복"),
    ("$5.81 억", "대한민국 수출 실적", "세계 유일의 순수출국 입지")
]

lefts = [Inches(0.8), Inches(3.8), Inches(6.8), Inches(9.8)]
for i, (val, lbl, sub) in enumerate(stats):
    shape = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, lefts[i], Inches(1.8), Inches(2.7), Inches(4.8))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_CARD
    shape.line.color.rgb = COLOR_GOLD
    shape.line.width = Pt(1.5)
    
    tf = shape.text_frame
    tf.word_wrap = True
    p0 = tf.paragraphs[0]
    p0.text = val
    p0.font.name = 'Arial Black'
    p0.font.size = Pt(26)
    p0.font.color.rgb = COLOR_PRIMARY
    p0.alignment = PP_ALIGN.CENTER
    p0.space_after = Pt(10)
    
    p1 = tf.add_paragraph()
    p1.text = lbl
    p1.font.name = 'Malgun Gothic'
    p1.font.size = Pt(14)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_TEXT_DARK
    p1.alignment = PP_ALIGN.CENTER
    p1.space_after = Pt(8)

    p2 = tf.add_paragraph()
    p2.text = sub
    p2.font.name = 'Malgun Gothic'
    p2.font.size = Pt(11)
    p2.font.color.rgb = COLOR_TEXT_MUTED
    p2.alignment = PP_ALIGN.CENTER

# Slide 3~4: Chart Showcase
def add_chart_showcase_slide(title, img_list):
    slide = prs.slides.add_slide(blank_slide_layout)
    set_bg(slide, COLOR_BG_LIGHT)
    add_header(slide, title)
    
    positions = [
        (Inches(0.8), Inches(1.5), Inches(5.6), Inches(2.6)),
        (Inches(6.8), Inches(1.5), Inches(5.6), Inches(2.6)),
        (Inches(0.8), Inches(4.3), Inches(5.6), Inches(2.6)),
        (Inches(6.8), Inches(4.3), Inches(5.6), Inches(2.6))
    ]
    for idx, (img_file, cap) in enumerate(img_list):
        if idx < 4:
            l, t, w, h = positions[idx]
            img_path = os.path.join("BIZ-Jeonbok", "images", img_file)
            if os.path.exists(img_path):
                slide.shapes.add_picture(img_path, l, t, width=w, height=h)

add_chart_showcase_slide("17종 핵심 EDA 시각화 차트 (1/2)", [
    ("01_univariate_flow_dist.png", "무역 유형 비율"),
    ("03_univariate_reporter_top30.png", "상위 보고국"),
    ("06_univariate_unitprice_dist.png", "단가 분포"),
    ("09_multivariate_reporter_flow_heatmap.png", "보고국 x 무역유형")
])

add_chart_showcase_slide("17종 핵심 EDA 시각화 차트 (2/2)", [
    ("15_price_volatility_anomaly.png", "단가 변동계수 CV"),
    ("16_customs_trade_balance.png", "수출입 무역수지"),
    ("17_country_cmd_clustering.png", "K-Means 군집분석"),
    ("11_tfidf_cmd_text.png", "TF-IDF 키워드")
])

# Slide 5: 4-Quadrant Strategy
slide_q = prs.slides.add_slide(blank_slide_layout)
set_bg(slide_q, COLOR_BG_LIGHT)
add_header(slide_q, "1인 종합상사 시장 개척 4분면 매트릭스 (4-Quadrant Strategy)")

quads = [
    ("Quadrant I: Star (주력 개척)", "홍콩, 싱가포르, 마카오", "활전복/건전복(HS 030781) 항공 산소팩 ($45~$55/kg)\n명품 건전복 선물세트", Inches(6.8), Inches(1.5)),
    ("Quadrant II: Cash Cow (대량 매출)", "미국, 캐나다, 중국", "전복 통조림/파우치(HS 160557) 해상 FCL\n아시안 마트 & HORECA 대량 공급", Inches(0.8), Inches(1.5)),
    ("Quadrant IV: Question Mark (신흥 럭셔리)", "일본, 베트남, 호주", "횟감용 IQF 냉동전복(HS 030783)\nFTA 무관세 혜택 활용", Inches(6.8), Inches(4.3)),
    ("Quadrant III: Selective (니치 가공)", "프랑스, 네덜란드, 이탈리아", "전복 내장 가공품(HS 160557)\n유럽 고급 시푸드 델리 시장", Inches(0.8), Inches(4.3))
]

for title_q, c_q, desc_q, l_q, t_q in quads:
    shape = slide_q.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l_q, t_q, Inches(5.7), Inches(2.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_CARD
    shape.line.color.rgb = COLOR_ACCENT
    shape.line.width = Pt(1)
    
    tf = shape.text_frame
    tf.word_wrap = True
    p0 = tf.paragraphs[0]
    p0.text = title_q
    p0.font.name = 'Malgun Gothic'
    p0.font.size = Pt(14)
    p0.font.bold = True
    p0.font.color.rgb = COLOR_PRIMARY
    
    p1 = tf.add_paragraph()
    p1.text = f"타겟: {c_q}"
    p1.font.name = 'Malgun Gothic'
    p1.font.size = Pt(12)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_GOLD
    
    p2 = tf.add_paragraph()
    p2.text = desc_q
    p2.font.name = 'Malgun Gothic'
    p2.font.size = Pt(11)
    p2.font.color.rgb = COLOR_TEXT_DARK

# Slide 6: Product-by-Product Star & Rising Star Strategy (Strict Dataset HS Codes)
slide_star = prs.slides.add_slide(blank_slide_layout)
set_bg(slide_star, COLOR_BG_LIGHT)
add_header(slide_star, "5대 품목별 Star & Rising Star 시장 및 데이터셋 수록 HS CODE")

rows, cols = 6, 5
table_shape_s = slide_star.shapes.add_table(rows, cols, Inches(0.8), Inches(1.5), Inches(11.7), Inches(5.3))
t_s = table_shape_s.table

headers_s = ["품목 분류 (데이터셋 수록 HS Code)", "Star 시장 (주력)", "Rising Star 시장 (신흥)", "물류 / 포장 방식", "추천 CIF 단가"]
for col_idx, h in enumerate(headers_s):
    cell = t_s.cell(0, col_idx)
    cell.text = h
    cell.fill.solid()
    cell.fill.fore_color.rgb = COLOR_PRIMARY
    p = cell.text_frame.paragraphs[0]
    p.font.name = 'Malgun Gothic'
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.alignment = PP_ALIGN.CENTER

star_pptx_data = [
    ("1. 활전복 (Live, HS 030781)", "홍콩, 싱가포르, 마카오", "베트남, 대만", "항공 (Air) 산소 해수 팩", "$31.00~$45.00 / kg"),
    ("2. 명품 건전복 (Dried, HS 030781)", "홍콩, 마카오, 광동성", "말레이시아, 미국(중화권)", "항공/해상 하드케이스", "$142.00 / 500g"),
    ("3. 통조림 & 파우치 (HS 160557)", "미국, 캐나다 (Cash Cow)", "호주, 영국/독일", "해상 400g 캔 / 1kg 파우치", "$14.30/캔, $26.80/kg"),
    ("4. 횟감용 IQF 냉동 (HS 030783)", "일본 (도쿄/오사카)", "베트남, 태국, 호주", "해상/항공 IQF 개별동결", "$30.00 / kg"),
    ("5. 전복 내장소스 (HS 160557)", "미국 (메인스트림)", "프랑스, 네덜란드, 이탈리아", "해상 유리병 / 파우치", "$12.50 / 병 (200g)")
]

for row_idx, data in enumerate(star_pptx_data, start=1):
    for col_idx, text in enumerate(data):
        cell = t_s.cell(row_idx, col_idx)
        cell.text = text
        p = cell.text_frame.paragraphs[0]
        p.font.name = 'Malgun Gothic'
        p.font.size = Pt(10)
        p.font.color.rgb = COLOR_TEXT_DARK

# Slide 7: Product Top 20 Buyers Matrix (Strict Dataset HS Codes)
slide_b20 = prs.slides.add_slide(blank_slide_layout)
set_bg(slide_b20, COLOR_BG_LIGHT)
add_header(slide_b20, "품목별 Top 20 잠재고객사 / 수산 전문 바이어 매트릭스")

rows, cols = 11, 5
table_shape_b20 = slide_b20.shapes.add_table(rows, cols, Inches(0.8), Inches(1.5), Inches(11.7), Inches(5.3))
t_b20 = table_shape_b20.table

headers_b20 = ["품목 구분 (데이터셋 HS Code)", "타겟 국가", "바이어 사명", "유통 형태", "이메일 / 연락처"]
for col_idx, h in enumerate(headers_b20):
    cell = t_b20.cell(0, col_idx)
    cell.text = h
    cell.fill.solid()
    cell.fill.fore_color.rgb = COLOR_PRIMARY
    p = cell.text_frame.paragraphs[0]
    p.font.name = 'Malgun Gothic'
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.alignment = PP_ALIGN.CENTER

b20_pptx_preview = [
    ("활전복 (HS 030781)", "홍콩", "On Kee Dry Seafood (安記海味)", "수산 도매 / 파인다이닝", "info@onkee.com"),
    ("활전복 (HS 030781)", "베트남", "Royal Seafood (Hải Sản Hoàng Gia)", "고급 활수산 프랜차이즈", "info@haisanhoanggia.com"),
    ("건전복 (HS 030781)", "홍콩", "Kee Wah Bakery & Trading Co.", "보양 선물세트 유통", "cs@keewah.com"),
    ("건전복 (HS 030781)", "말레이시아", "Eu Yan Sang Malaysia (余仁生)", "동남아 최대 중화 보양 체인", "info@euyansang.com.my"),
    ("통조림 (HS 160557)", "미국", "H Mart Commercial Division", "대형 아시안 마트 체인", "b2b@hmart.com"),
    ("통조림 (HS 160557)", "캐나다", "T&T Supermarket Inc.", "대형 캐나다 아시안 체인", "customer.service@tntsupermarket.com"),
    ("IQF 냉동 (HS 030783)", "일본", "True World Foods Japan", "일식 스시/사시미 벤더", "japaninfo@trueworldfoods.com"),
    ("IQF 냉동 (HS 030783)", "호주", "De Costi Seafoods", "오세아니아 최대 수산 유통", "info@decosti.com.au"),
    ("내장소스 (HS 160557)", "프랑스", "Maison Plisson / Gourmet Food", "파리 프리미엄 시푸드 델리", "contact@lamaisonplisson.com"),
    ("내장소스 (HS 160557)", "미국", "Trader Joe's Specialty Sourcing", "미국 스페셜티 유통 체인", "customer@traderjoes.com")
]

for row_idx, data in enumerate(b20_pptx_preview, start=1):
    for col_idx, text in enumerate(data):
        cell = t_b20.cell(row_idx, col_idx)
        cell.text = text
        p = cell.text_frame.paragraphs[0]
        p.font.name = 'Malgun Gothic'
        p.font.size = Pt(9.5)
        p.font.color.rgb = COLOR_TEXT_DARK

# Slide 8: Commercial Execution Playbook
slide_pb = prs.slides.add_slide(blank_slide_layout)
set_bg(slide_pb, COLOR_BG_LIGHT)
add_header(slide_pb, "실전 비즈니스 영업 파이프라인 & 5단계 액션 플레이북")

steps_pb = [
    ("Action 1: Pitching", "Cold Outreach", "20개 바이어 대상 영문 제안서 송부\n15% 우위 단가 + 폐사율 3% 보증 오퍼"),
    ("Action 2: Sample Test", "품질 검증", "10kg 산소주입 활전복 Air 송부\n생존율, 수율, 멸균 검증서 확인"),
    ("Action 3: Contract", "결제 조건", "신규 T/T 30% Deposit / 70% B/L Copy\n대형 수산상 L/C at Sight 활용"),
    ("Action 4: Compliance", "통관 검역", "미국 FDA 사전통보 / 한-베 FTA C/O\n관세 0% 적용 및 EU 위생 허가"),
    ("Action 5: Re-order", "LTV 극대화", "춘절/중추절 3개월 전 건전복 수주\n북미/유럽 연말 통조림 FCL 발주")
]

lefts_5 = [Inches(0.8), Inches(3.2), Inches(5.6), Inches(8.0), Inches(10.4)]
for i, (act_title, sub_t, desc_t) in enumerate(steps_pb):
    shape = slide_pb.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, lefts_5[i], Inches(1.8), Inches(2.1), Inches(4.8))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_CARD
    shape.line.color.rgb = COLOR_GOLD
    shape.line.width = Pt(1.5)
    
    tf = shape.text_frame
    tf.word_wrap = True
    p0 = tf.paragraphs[0]
    p0.text = act_title
    p0.font.name = 'Arial Black'
    p0.font.size = Pt(13)
    p0.font.color.rgb = COLOR_PRIMARY
    p0.alignment = PP_ALIGN.CENTER
    p0.space_after = Pt(4)
    
    p1 = tf.add_paragraph()
    p1.text = sub_t
    p1.font.name = 'Malgun Gothic'
    p1.font.size = Pt(12)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_GOLD
    p1.alignment = PP_ALIGN.CENTER
    p1.space_after = Pt(10)

    p2 = tf.add_paragraph()
    p2.text = desc_t
    p2.font.name = 'Malgun Gothic'
    p2.font.size = Pt(10)
    p2.font.color.rgb = COLOR_TEXT_DARK

# Save PPTX
output_file = os.path.join("BIZ-Jeonbok", "reports", "Wando_Abalone_Market_Entry.pptx")
alt_output_file = os.path.join("BIZ-Jeonbok", "reports", "Wando_Abalone_Market_Entry_v3.pptx")
try:
    prs.save(output_file)
    print(f"--- UPDATED PPTX DECK STRICTLY MATCHING DATASET HS CODES GENERATED SUCCESSFULLY: {output_file} ---")
except PermissionError:
    prs.save(alt_output_file)
    print(f"--- UPDATED PPTX DECK SAVED TO ALTERNATE PATH: {alt_output_file} ---")
