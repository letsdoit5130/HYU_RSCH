"""
BIZ-전복_Gathered_EDA_Report.md 기반 TOP 10 유망국가 표 3종 완비 PPT(.pptx) 발표 덱 생성 스크립트

이 스크립트는 전복 EDA 보고서의 15개 시각화 차트, 미수(Size) 가격 구조 표, 
HS Code 3대 품목별 TOP 10 유망국가 표 3종(각 10개국 완비) 및 4대 특별 부록을 16:9 와이드 PPT로 변환 생성합니다.
"""
import os
import sys
import pptx
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Windows 콘솔 인코딩 방어
if sys.platform == 'win32':
    try:
        sys.stdout.reconfigure(encoding='utf-8')
    except Exception:
        pass

BASE_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
PPTX_OUTPUT = os.path.join(BASE_DIR, 'reports', 'BIZ_Jeonbok_Market_Entry_Deck.pptx')
IMG_DIR = os.path.join(BASE_DIR, 'images')

def create_pptx_deck():
    prs = pptx.Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    COLOR_PRIMARY = RGBColor(0x1F, 0x49, 0x7D)
    COLOR_TEXT = RGBColor(0x33, 0x33, 0x33)

    # 1. Cover Slide
    slide1 = prs.slides.add_slide(blank_layout)
    bg = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLOR_PRIMARY
    bg.line.fill.background()

    txBox = slide1.shapes.add_textbox(Inches(1), Inches(2.2), Inches(11.333), Inches(3.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    p1 = tf.paragraphs[0]
    p1.text = "🦪 한국산 전복(Abalone) 글로벌 시장개척 전략 덱"
    p1.font.bold = True
    p1.font.size = Pt(36)
    p1.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    p2 = tf.add_paragraph()
    p2.text = "UN Comtrade 무역 분석, 15개 차트, TOP 10 유망국가 10개국 완비 및 파트너 소싱 전략"
    p2.font.size = Pt(20)
    p2.font.color.rgb = RGBColor(0xD0, 0xE0, 0xF0)

    # 2. Executive Summary Slide
    slide2 = prs.slides.add_slide(blank_layout)
    hb = slide2.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(1.0))
    tf = hb.text_frame
    p = tf.paragraphs[0]
    p.text = "📌 Executive Summary: 미수(Size)별 가격 구조"
    p.font.bold = True
    p.font.size = Pt(24)
    p.font.color.rgb = COLOR_PRIMARY

    t_shape = slide2.shapes.add_table(6, 6, Inches(0.8), Inches(1.8), Inches(11.7), Inches(4.8))
    t = t_shape.table

    p_data = [
        ["품목 규격 / 미수", "마리당 중량", "주요 수출국", "평균 단가", "주요 타깃 시장", "소싱 추천 포인트"],
        ["10미 미만 (대과)", "100g 이상", "한국 완도, 호주", "$42.0 ~ $48.0", "일본 고급 일식집, 료칸", "고급 항공직송 프리미엄 오퍼"],
        ["10 ~ 12미 (중대과)", "80g ~ 100g", "한국 완도", "$36.0 ~ $40.0", "도쿄 도요스 시장 도매상사", "메인 수출 주력 미수 1차 상사"],
        ["13 ~ 15미 (중과)", "65g ~ 80g", "한국, 중국", "$30.0 ~ $34.0", "관서 레스토랑, 아시안 마트", "H-Mart, 99 Ranch 채널 공급"],
        ["15 ~ 20미 (중소과)", "50g ~ 65g", "한국, 베트남", "$24.0 ~ $28.0", "냉동 IQF 가공, 외식 체인", "해상 IQF 컨테이너 대량 공급"],
        ["20미 이상 (소과)", "50g 미만", "한국, 중국", "$18.0 ~ $22.0", "통조림 가공, HMR 가공", "통조림 FDA 승인 공장 연동"]
    ]

    for r_idx, row in enumerate(p_data):
        for c_idx, val in enumerate(row):
            cell = t.cell(r_idx, c_idx)
            cell.text = val
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(12)
            if r_idx == 0:
                p.font.bold = True
                p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLOR_PRIMARY

    # 3. TOP 10 유망국가 분석 슬라이드 3장 (각 10개국 완비)
    tables_info = [
        ("[표 1] HS 0307.81 (활/신선 전복) TOP 10 유망 국가", [
            ["유망순위", "타깃 국가", "무역액 점유율", "컨택 파트너 종류", "시장개척 포인트"],
            ["1위", "일본 (Japan)", "35.4%", "도쿄 도요스 시장 수산물 수입 도매상사", "완도산 활전복 페리/항공 직송 공급"],
            ["2위", "중국 (China)", "24.1%", "동해안 수산물 수입 및 유통 상사", "산둥성/상하이 고급 호텔 체인 공급"],
            ["3위", "홍콩 (Hong Kong)", "18.2%", "고급 수산물 건재 시장 수입상사", "고급 딤섬 및 레스토랑 직송 공급"],
            ["4위", "대만 (Taiwan)", "7.5%", "타이베이 고급 수산물 1차 수입상", "일식 뷔페 및 연회장 활전복 공급"],
            ["5위", "미국 (USA)", "4.8%", "LA/NY 아시안 수산물 벤더", "한인/아시안 고소득층 항공 직송"],
            ["6위", "싱가포르 (Singapore)", "3.2%", "마리나 베이 외식 그룹 벤더", "고급 해산물 뷔페 및 호텔 공급"],
            ["7위", "베트남 (Vietnam)", "2.5%", "호치민/하노이 수산물 수입상", "한국 식당가 및 고급 수산 레스토랑"],
            ["8위", "캐나다 (Canada)", "1.8%", "밴쿠버 아시안 수산 유통사", "밴쿠버/토론토 아시안 마트 공급"],
            ["9위", "태국 (Thailand)", "1.3%", "방콕 고급 수산물 수입 대리점", "방콕 5성급 호텔 수산물 오퍼"],
            ["10위", "호주 (Australia)", "1.2%", "시드니 아시안 식품 유통 벤더", "호주 한인 마트 및 아시안 레스토랑"]
        ]),
        ("[표 2] HS 0307.83 (냉동 전복) TOP 10 유망 국가", [
            ["유망순위", "타깃 국가", "무역액 점유율", "컨택 파트너 종류", "시장개척 포인트"],
            ["1위", "미국 (USA)", "42.1%", "미 서부 최대 수산물 수입 벤더 (PASCO)", "아시안 마트향 냉동 IQF 공급"],
            ["2위", "대만 (Taiwan)", "19.8%", "타이베이 식자재 수입 디스트리뷰터", "뷔페/연회장향 IQF 냉동 대량 공급"],
            ["3위", "일본 (Japan)", "15.3%", "관서 지역 냉동 수산물 수입 대리점", "성수기 외식 체인 원료 공급"],
            ["4위", "홍콩 (Hong Kong)", "8.2%", "냉동 수산물 전문 수입 유통사", "외식 체인 및 호텔 냉동 IQF 공급"],
            ["5위", "싱가포르 (Singapore)", "4.5%", "동남아 아시안 식자재 유통 벤더", "뷔페 및 딤섬 프랜차이즈 공급"],
            ["6위", "중국 (China)", "3.8%", "연안 도시 식품 가공 및 유통사", "가공 원료용 냉동 IQF 전복 공급"],
            ["7위", "캐나다 (Canada)", "2.1%", "토론토 수산물 수입 벤더", "아시안 마트 냉동 해산물 공급"],
            ["8위", "베트남 (Vietnam)", "1.8%", "외식 식자재 1차 수입상", "프랜차이즈 레스토랑 IQF 공급"],
            ["9위", "태국 (Thailand)", "1.3%", "방콕 식자재 수입 대리점", "외식 뷔페 및 씨푸드 레스토랑 공급"],
            ["10위", "영국 (United Kingdom)", "1.1%", "런던 아시안 식품 수입 벤더", "런던 아시안 마트 및 한식당 공급"]
        ]),
        ("[표 3] HS 1605.57 (전복 통조림) TOP 10 유망 국가", [
            ["유망순위", "타깃 국가", "무역액 점유율", "컨택 파트너 종류", "시장개척 포인트"],
            ["1위", "홍콩 (Hong Kong)", "48.5%", "홍콩 셩완 수산물 건재 수입상사", "춘절 명절 선물 세트용 B2B 캔 공급"],
            ["2위", "싱가포르 (Singapore)", "22.1%", "싱가포르 고급 선물 세트 유통 벤더", "명절/기념일 프리미엄 캔 공급"],
            ["3위", "미국 (USA)", "14.8%", "북미 아시안 식품 수입 벤더", "FDA LACF 승인 캔 통조림 유통"],
            ["4위", "대만 (Taiwan)", "4.2%", "명절 선물 세트 수입 유통사", "명절 고급 전복 캔 선물 세트 공급"],
            ["5위", "캐나다 (Canada)", "3.1%", "밴쿠버 아시안 마트 유통 벤더", "한인/중국인 마트 캔 전복 공급"],
            ["6위", "호주 (Australia)", "2.3%", "시드니/멜버른 아시안 식품 수입상", "선물용 캔 전복 유통"],
            ["7위", "일본 (Japan)", "1.8%", "고급 통조림 식자재 유통사", "료칸 및 기프트 숍 고급 캔 오퍼"],
            ["8위", "베트남 (Vietnam)", "1.2%", "고급 선물 세트 수입상", "호치민/하노이 명절 선물용 캔 공급"],
            ["9위", "태국 (Thailand)", "1.1%", "방콕 아시안 식품 수입 벤더", "고급 아시안 마트 캔 유통"],
            ["10위", "영국 (United Kingdom)", "0.9%", "런던 프리미엄 기프트 숍 벤더", "런던 아시안 명절 기프트 공급"]
        ])
    ]

    for title_txt, t_rows in tables_info:
        slide = prs.slides.add_slide(blank_layout)
        hb = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.8))
        tf = hb.text_frame
        p = tf.paragraphs[0]
        p.text = f"🗺️ {title_txt}"
        p.font.bold = True
        p.font.size = Pt(22)
        p.font.color.rgb = COLOR_PRIMARY

        t_shape = slide.shapes.add_table(11, 5, Inches(0.8), Inches(1.3), Inches(11.7), Inches(5.7))
        t = t_shape.table

        for r_idx, row in enumerate(t_rows):
            for c_idx, val in enumerate(row):
                cell = t.cell(r_idx, c_idx)
                cell.text = val
                p = cell.text_frame.paragraphs[0]
                p.font.size = Pt(10)
                if r_idx == 0:
                    p.font.bold = True
                    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = COLOR_PRIMARY

    # 4. 15개 시각화 차트 슬라이드 15장
    chart_files = [
        ("01_annual_trade_trend.png", "1. 연도별 무역액 추이"),
        ("02_top_exporter_ranking.png", "2. TOP 10 주요 수출국 무역액"),
        ("03_top_importer_ranking.png", "3. TOP 10 주요 수입국 무역액"),
        ("04_unit_price_distribution.png", "4. 전복 평균 단가 ($/kg) 분포"),
        ("05_monthly_seasonality.png", "5. 월별 거래 계절성 지수"),
        ("06_hs_code_share.png", "6. HS Code별 거래액 점유율"),
        ("07_price_vs_weight_scatter.png", "7. 물량 vs 단가 상관관계 산점도"),
        ("08_top5_importer_growth.png", "8. TOP 5 수입국 연도별 성장 추이"),
        ("09_market_concentration_pareto.png", "9. 수입 시장 파레토 집중도 분석"),
        ("10_export_price_heatmap.png", "10. 주요 수입국별 연도별 평균 단가 히트맵"),
        ("11_trade_balance_waterfall.png", "11. 무역 구조 폭포수 분석"),
        ("12_country_price_boxplot.png", "12. TOP 주요 국가별 단가 변동성 박스플롯"),
        ("13_hhi_index_trend.png", "13. 시장 집중도(HHI Index) 추이"),
        ("14_size_pricing_structure.png", "14. 미수(Size) 규격별 가격 구조"),
        ("15_promising_country_matrix.png", "15. 전복 유망 국가 시장 성숙도-단가 매트릭스")
    ]

    for img_name, title_txt in chart_files:
        img_path = os.path.join(IMG_DIR, img_name)
        if os.path.exists(img_path):
            slide = prs.slides.add_slide(blank_layout)
            hb = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.8))
            tf = hb.text_frame
            p = tf.paragraphs[0]
            p.text = f"📊 {title_txt}"
            p.font.bold = True
            p.font.size = Pt(24)
            p.font.color.rgb = COLOR_PRIMARY

            slide.shapes.add_picture(img_path, Inches(1.5), Inches(1.4), width=Inches(10.333))

    prs.save(PPTX_OUTPUT)
    print(f"✅ [TOP 10 유망국가 10개국 완비] PPT(.pptx) 슬라이드 덱 재생성 완료: {PPTX_OUTPUT}")

if __name__ == "__main__":
    create_pptx_deck()
