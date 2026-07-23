"""
이 스크립트는 예스24 베스트셀러 고도화 EDA 보고서 데이터를 바탕으로,
30페이지 분량의 실제 파워포인트(.pptx) 발표 자료를 자동 생성해 주는 프로그램입니다.
주요 기능:
- 제목 폰트("G마켓 산스 Bold"), 내용 폰트("나눔 고딕")의 정밀 서식 바인딩
- 모든 슬라이드에 인포그래픽형 카드 도형(ROUNDED_RECTANGLE) 및 비교 그리드 자동 생성
- 13개의 실제 데이터 분석 차트 이미지(yes24/images)를 2열 비대칭 구조로 정상 임베딩
- 슬라이드 하단에 발표 시 바로 읽을 수 있는 2분 분량(약 450자)의 한국어 발표자 노트(Speaker Note) 자동 추가
- 16:9 와이드스크린 화면비 고정 및 앤틱 다크아카데미아/골드/차콜 배색 테마 적용
"""
# -*- coding: utf-8 -*-
import os
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

# 1. 파일 경로 정의
base_dir = r"C:\Users\leeak\OneDrive\1.HaeYu\HYU_RSCH\yes24"
image_dir = os.path.join(base_dir, "images")
output_path = os.path.join(base_dir, "docs", "bestsellers_presentation.pptx")

# 2. 색상 상수 정의 (다크아카데미아 & 앤틱 골드)
BG_DARK = RGBColor(35, 31, 32)      # 딥 차콜 브라운 #231F20 (간지 및 표지용)
BG_LIGHT = RGBColor(248, 250, 248)  # 소프트 오프화이트 #F8FAF8 (본문용)
TXT_DARK = RGBColor(17, 17, 17)      # near-black #111111
TXT_LIGHT = RGBColor(230, 213, 195)  # 크림 베이지 #E6D5C3
COLOR_GOLD = RGBColor(191, 163, 138) # 앤틱 골드 #BFA38A (강조색)
COLOR_CARD_BG = RGBColor(255, 255, 255) # 카드 배경 (흰색)
COLOR_BORDER = RGBColor(220, 210, 200)  # 카드 테두리

def set_font(run, font_name, size_pt, bold=False, italic=False, color_rgb=None):
    """텍스트 런(Run)에 폰트 이름, 크기, 굵기, 기울임 및 색상을 지정하는 헬퍼 함수"""
    run.font.name = font_name
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.italic = italic
    if color_rgb:
        run.font.color.rgb = color_rgb

def set_slide_background(slide, color_rgb):
    """슬라이드의 단색 배경색을 지정하는 함수"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color_rgb

def add_card_shape(slide, left, top, width, height, fill_rgb=COLOR_CARD_BG, border_rgb=COLOR_BORDER):
    """둥근 직사각형 인포그래픽 카드를 생성하고 색상과 테두리를 설정하는 함수"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    shape.line.color.rgb = border_rgb
    shape.line.width = Pt(1.5)
    return shape

def add_title(slide, text, is_dark_bg=False):
    """슬라이드 제목을 추가하고 폰트(G마켓 산스 Bold)를 설정하는 함수"""
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12.0), Inches(0.8))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    set_font(run, "Gmarket Sans", 32, bold=True, color_rgb=TXT_LIGHT if is_dark_bg else TXT_DARK)

def build_presentation():
    prs = Presentation()
    
    # 16:9 와이드스크린 비율 설정
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # 빈 슬라이드 레이아웃(6번 레이아웃) 사용
    blank_layout = prs.slide_layouts[6]
    
    # -------------------------------------------------------------
    # Slide 1. 표지 (Title Slide)
    # -------------------------------------------------------------
    slide1 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide1, BG_DARK)
    
    # 제목 카드 데코레이션
    add_card_shape(slide1, Inches(0.8), Inches(1.8), Inches(11.73), Inches(3.8), fill_rgb=RGBColor(45, 40, 41), border_rgb=COLOR_GOLD)
    
    title_box = slide1.shapes.add_textbox(Inches(1.2), Inches(2.2), Inches(11.0), Inches(2.0))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "YES24 베스트셀러 고도화 EDA 종합 보고서"
    set_font(p.runs[0], "Gmarket Sans", 42, bold=True, color_rgb=TXT_LIGHT)
    
    p2 = tf.add_paragraph()
    p2.text = "데이터 기반 도서 시장의 구조 진단 및 출판 기획 전략 제언"
    set_font(p2.runs[0], "Nanum Gothic", 18, color_rgb=COLOR_GOLD)
    
    p3 = tf.add_paragraph()
    p3.text = "\n발표자: 전문 데이터 분석가  |  날짜: 2026년 7월 13일"
    set_font(p3.runs[0], "Nanum Gothic", 14, color_rgb=TXT_LIGHT)
    
    note1 = (
        "안녕하십니까, 오늘 발표를 맡은 전문 데이터 분석가입니다. 오늘 여러분께 보고드릴 자료는 "
        "국내 대표 도서 플랫폼인 예스24의 실제 베스트셀러 데이터를 활용해 정량적인 관점에서 도서 흥행의 "
        "구조와 독자들의 구매 특성을 입증한 고도화 EDA 결과물입니다. 직관과 트렌드 감각에만 의존하던 기존의 "
        "출판 기획 관행에서 탈피하여, 판매지수, 평점, 할인율, 리뷰 건수, 그리고 태그 텍스트 데이터 마이닝을 "
        "아우르는 30대 핵심 지표 분석과 이에 기반한 성공 프레임워크를 공유하겠습니다. 총 2분 분량의 장표별 설명을 "
        "참고해 주시며 비즈니스적 통찰력을 발굴해 보시기 바랍니다."
    )
    slide1.notes_slide.notes_text_frame.text = note1

    # -------------------------------------------------------------
    # Slide 2. 목차 (Agenda Slide)
    # -------------------------------------------------------------
    slide2 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide2, BG_LIGHT)
    add_title(slide2, "발표 목차 (Agenda)")
    
    # 4개 아젠다 영역 카드 배치
    agendas = [
        ("PART 1. 데이터 품질 진단", "데이터 무결성, info() 형태 분석 및 전처리 절차 점검"),
        ("PART 2. 단일 변수 통계 분석", "출판사, 저자 빈도 및 가격, 평점, 판매지수의 밀집도 진단"),
        ("PART 3. 이변량 교차 영향성", "출판사별 성과 편차 및 평점 만족도-판매지수 디커플링 실증"),
        ("PART 4. 다변량 및 텍스트 마이닝", "태그 TF-IDF 자연어 가중치 추출 및 리뷰 플라이휠 효과 규명")
    ]
    
    for i, (ag_t, ag_d) in enumerate(agendas):
        x = Inches(0.8 + (i * 2.95))
        add_card_shape(slide2, x, Inches(1.8), Inches(2.8), Inches(4.5))
        tb = slide2.shapes.add_textbox(x + Inches(0.1), Inches(2.0), Inches(2.6), Inches(4.1))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"0{i+1}\n\n{ag_t}"
        set_font(p.runs[0], "Gmarket Sans", 18, bold=True, color_rgb=TXT_DARK)
        p2 = tf.add_paragraph()
        p2.text = f"\n{ag_d}"
        set_font(p2.runs[0], "Nanum Gothic", 13, color_rgb=RGBColor(80, 80, 80))
        
    note2 = (
        "전체 목차는 총 4대 분석 도메인과 이후에 이어질 복합 상관관계 분석, 종합 인사이트 도출, 전략적 제언 "
        "파트로 세분되어 구성되어 있습니다. 파트 1에서는 로우데이터의 877행 15열 형태를 뜯어보고 중복 여부 및 "
        "결측치 전처리 이력을 말씀드립니다. 파트 2에서는 일변량 기술통계를 통해 시장의 원초적 기초 체력을 측정하고, "
        "파트 3에서는 본격적인 연관성 분석을 위해 산점도를 통해 통계적 요인들을 입증합니다. 마지막 파트 4에서는 "
        "TF-IDF 자연어 키워드와 리뷰 건수의 임계점 분석을 통해 실질적인 성공 프레임워크를 유도해 드리겠습니다."
    )
    slide2.notes_slide.notes_text_frame.text = note2

    # -------------------------------------------------------------
    # Slide 3. [간지] PART 1. 데이터 기본 정보 및 품질 파악
    # -------------------------------------------------------------
    slide3 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide3, BG_DARK)
    
    add_card_shape(slide3, Inches(1.0), Inches(2.0), Inches(11.33), Inches(3.5), fill_rgb=RGBColor(40, 35, 36), border_rgb=COLOR_GOLD)
    tb = slide3.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(10.33), Inches(2.5))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "PART 1. 데이터 기본 정보 및 품질 파악"
    set_font(p.runs[0], "Gmarket Sans", 36, bold=True, color_rgb=TXT_LIGHT)
    p2 = tf.add_paragraph()
    p2.text = "예스24 베스트셀러 데이터셋의 정제 상태와 기본 통계량 진단"
    set_font(p2.runs[0], "Nanum Gothic", 18, color_rgb=COLOR_GOLD)
    
    note3 = (
        "첫 번째 대분류인 데이터 기본 정보 및 품질 파악 파트입니다. 아무리 훌륭한 시각화나 통계 분석 기법을 "
        "적용하더라도 그 재료가 되는 데이터셋에 누락, 왜곡, 혹은 타입 상의 기계적 오류가 있다면 분석 결과의 "
        "신뢰도가 통째로 흔들리게 됩니다. 따라서 본 단락에서는 877개 도서 데이터셋의 세부 스키마 형태와 전수 검증 "
        "결과를 상세하게 검증해 보며, 데이터 분석가로서 무결성을 입증받는 정제 로직을 공유하겠습니다."
    )
    slide3.notes_slide.notes_text_frame.text = note3

    # -------------------------------------------------------------
    # Slide 4. 데이터 무결성 및 개요
    # -------------------------------------------------------------
    slide4 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide4, BG_LIGHT)
    add_title(slide4, "1.1 분석 대상 데이터 개요 및 차원")
    
    # 3대 KPI 카드 배치
    kpis = [
        ("전체 수집 데이터셋 차원", "877 행 X 15 열", "도서명, 저자, 출판사, 판매지수, 평점, 태그 등 15개 고유 특징 변수 확보"),
        ("데이터 무결성 검증", "중복 데이터 0 건", "전수 행 중복 여부(duplicated) 검증 결과 무결한 고유 도서들로 판명"),
        ("핵심 수치형/범주형 분류", "수치 7대 / 범주 8대", "상품번호, 순위, 정가, 할인가, 판매지수, 리뷰건수, 평점 외 텍스트 메타데이터 분류")
    ]
    for i, (k_t, k_n, k_d) in enumerate(kpis):
        x = Inches(0.8 + (i * 3.95))
        add_card_shape(slide4, x, Inches(2.0), Inches(3.8), Inches(4.0))
        tb = slide4.shapes.add_textbox(x + Inches(0.2), Inches(2.2), Inches(3.4), Inches(3.6))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = k_t
        set_font(p.runs[0], "Nanum Gothic", 15, bold=True, color_rgb=COLOR_GOLD)
        p2 = tf.add_paragraph()
        p2.text = f"\n{k_n}"
        set_font(p2.runs[0], "Gmarket Sans", 24, bold=True, color_rgb=TXT_DARK)
        p3 = tf.add_paragraph()
        p3.text = f"\n{k_d}"
        set_font(p3.runs[0], "Nanum Gothic", 12, color_rgb=RGBColor(100, 100, 100))
        
    note4 = (
        "데이터셋의 전체 스케일과 무결성 수치입니다. 수집된 도서는 총 877권이며 세부 열 속성은 15개로 구성됩니다. "
        "중복 적재 오류를 검증하기 위해 Pandas duplicated() 로직을 가동해 전수 식별자 조사를 진행한 결과, 단 한 건의 "
        "중복 행도 존재하지 않는 순수 고유값 도서군임을 증명해 보였습니다. 이는 우리가 앞으로 전개할 일변량, 이변량 "
        "상관 분석의 백데이터로서 신뢰도 높은 토대를 제공함을 약속합니다."
    )
    slide4.notes_slide.notes_text_frame.text = note4

    # -------------------------------------------------------------
    # Slide 5. 상위 5개행 도서 목록 (테이블 인포그래픽)
    # -------------------------------------------------------------
    slide5 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide5, BG_LIGHT)
    add_title(slide5, "1.2 상위 베스트셀러 5개행 샘플 데이터")
    
    # 5개 데이터 카드 가로 행 배치
    top_samples = [
        ("1위", "나의 첫 번째 부동산 교과서", "송희구 / 서삼독", "168,537점", "9.3점"),
        ("2위", "부의 갈림길", "오건영 / 포레스트북스", "155,631점", "9.9점"),
        ("3위", "주식 투자를 잘한다는 것", "육과장 / 노티스", "109,503점", "9.8점"),
        ("4위", "독하게 돈 공부", "박소연 / 메이븐", "4,878점", "8.0점"),
        ("5위", "박곰희 연금 부자 수업", "박곰희 / 인플루엔셜", "563,742점", "9.7점")
    ]
    for i, (rank, title, info, sales, rating) in enumerate(top_samples):
        y = Inches(1.7 + (i * 1.05))
        add_card_shape(slide5, Inches(0.8), y, Inches(11.73), Inches(0.95))
        
        tb = slide5.shapes.add_textbox(Inches(1.0), y + Inches(0.1), Inches(11.33), Inches(0.75))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"[{rank}]  {title}  |  저자/출판사: {info}  |  판매지수: {sales}  |  평점: {rating}"
        set_font(p.runs[0], "Nanum Gothic", 14, bold=True, color_rgb=TXT_DARK)
        
    note5 = (
        "최상위 1위부터 5위까지의 도서 목록입니다. 1위부터 5위의 도서명 구성을 직관적으로 살펴보면, 송희구의 "
        "부동산 교과서, 오건영의 부의 갈림길, 박곰희의 연금 부자 수업 등 전형적인 자산 증식과 경제적 해법을 던지는 "
        "실용 재테크 서적군이 시장 최상단을 독차지하고 있습니다. 여기서 통계적 특이점은 5위인 박곰희 연금 부자 수업이 "
        "56만 점의 폭발적 판매지수를 내며 1위보다 큰 판매 볼륨을 내고 있다는 것이며, 4위 도서는 4,800점대로 급감해 "
        "순위와 실제 판매 스코어가 비선형으로 단절되어 있음을 포착할 수 있습니다."
    )
    slide5.notes_slide.notes_text_frame.text = note5

    # -------------------------------------------------------------
    # Slide 6. 하위 5개행 도서 목록 (테이블 인포그래픽)
    # -------------------------------------------------------------
    slide6 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide6, BG_LIGHT)
    add_title(slide6, "1.3 하위 베스트셀러 5개행 샘플 데이터")
    
    bot_samples = [
        ("873위", "소유의 종말", "제러미 리프킨 / 민음사", "1,452점", "8.8점"),
        ("874위", "인사이더 인사이트", "이용준 / 에프엔미디어", "30,120점", "9.5점"),
        ("875위", "국제경제론", "김신행, 김태기 / 법문사", "726점", "0.0점"),
        ("876위", "엘리어트 파동이론 마스터", "글렌 닐리 / 원앤원북스", "4,272점", "7.4점"),
        ("877위", "영업의 神신 100법칙", "하야카와 마사루 / 지상사", "2,553점", "9.9점")
    ]
    for i, (rank, title, info, sales, rating) in enumerate(bot_samples):
        y = Inches(1.7 + (i * 1.05))
        add_card_shape(slide6, Inches(0.8), y, Inches(11.73), Inches(0.95))
        
        tb = slide6.shapes.add_textbox(Inches(1.0), y + Inches(0.1), Inches(11.33), Inches(0.75))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"[{rank}]  {title}  |  저자/출판사: {info}  |  판매지수: {sales}  |  평점: {rating}"
        set_font(p.runs[0], "Nanum Gothic", 14, bold=True, color_rgb=TXT_DARK)
        
    note6 = (
        "최하위권 873위부터 877위까지의 데이터 구성입니다. 민음사의 소유의 종말과 같은 고전 명작이 눈에 띄며, "
        "김신행의 국제경제론이나 엘리어트 파동이론처럼 대학 교재 및 실무 전공 서적이 위치하고 있습니다. 하위권 도서의 "
        "판매지수는 주로 1,000점 내외이거나 그 미만인 726점 수준으로 저조하게 적재되어 있으며, 평점의 경우에도 "
        "0.0점(리뷰 전무) 혹은 7.4점 등 상위권 대비 엇갈린 고객 평가가 관찰되어 베스트셀러 내에서도 쏠림 현상이 "
        "분명하게 드러나기 시작합니다."
    )
    slide6.notes_slide.notes_text_frame.text = note6

    # -------------------------------------------------------------
    # Slide 7. [간지] PART 2. 일변량 기술통계 및 단일 변수 분석
    # -------------------------------------------------------------
    slide7 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide7, BG_DARK)
    
    add_card_shape(slide7, Inches(1.0), Inches(2.0), Inches(11.33), Inches(3.5), fill_rgb=RGBColor(40, 35, 36), border_rgb=COLOR_GOLD)
    tb = slide7.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(10.33), Inches(2.5))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "PART 2. 일변량 기술통계 및 단일 변수 분석"
    set_font(p.runs[0], "Gmarket Sans", 36, bold=True, color_rgb=TXT_LIGHT)
    p2 = tf.add_paragraph()
    p2.text = "도서 개별 정보, 가격, 판매량의 기초 체력 및 데이터 밀집 영역 검사"
    set_font(p2.runs[0], "Nanum Gothic", 18, color_rgb=COLOR_GOLD)
    
    note7 = (
        "두 번째 단분류 영역인 일변량 기술통계 및 단일 변수 분석 파트입니다. 이 단락에서는 도서 개별 변수들의 "
        "기초 통계 수치를 기술하고, 출판사 및 저자의 베스트셀러 등록 밀집 빈도를 낱낱이 파악해 보겠습니다. 더불어, "
        "도서의 가격 포지셔닝 분포와 평점의 편향 형태, 판매지수의 불균형 상태를 가로 막대와 히스토그램 차트 "
        "이미지 5개를 우측과 좌측에 차례로 임베딩하여 시각적이고 전문적으로 입증하겠습니다."
    )
    slide7.notes_slide.notes_text_frame.text = note7

    # -------------------------------------------------------------
    # Slide 8. 수치형/범주형 기초 통계 요약 (인포그래픽 카드)
    # -------------------------------------------------------------
    slide8 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide8, BG_LIGHT)
    add_title(slide8, "2.1 수치형 및 범주형 데이터 변수 요약")
    
    # 2열 비교 카드 배치
    add_card_shape(slide8, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8))
    tb_l = slide8.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(5.2), Inches(4.4))
    tf_l = tb_l.text_frame
    tf_l.word_wrap = True
    p_l = tf_l.paragraphs[0]
    p_l.text = "수치형 변수 기초 요약"
    set_font(p_l.runs[0], "Gmarket Sans", 20, bold=True, color_rgb=TXT_DARK)
    p_l2 = tf_l.add_paragraph()
    p_l2.text = (
        "\n■ 평균 정가: 21,830원 (중앙값 21,000원)\n"
        "■ 평균 할인율: 9.9% (도서정가제 10%의 고착화)\n"
        "■ 평균 판매지수: 12,219.9점 (중앙값 3,546.0점)\n"
        "■ 평균 평점: 8.6점 (중앙값 9.5점의 극심한 고평가 쏠림)"
    )
    set_font(p_l2.runs[0], "Nanum Gothic", 14, color_rgb=RGBColor(60, 60, 60))
    
    add_card_shape(slide8, Inches(6.8), Inches(1.8), Inches(5.6), Inches(4.8))
    tb_r = slide8.shapes.add_textbox(Inches(7.0), Inches(2.0), Inches(5.2), Inches(4.4))
    tf_r = tb_r.text_frame
    tf_r.word_wrap = True
    p_r = tf_r.paragraphs[0]
    p_r.text = "범주형 변수 기초 요약"
    set_font(p_r.runs[0], "Gmarket Sans", 20, bold=True, color_rgb=TXT_DARK)
    p_r2 = tf_r.add_paragraph()
    p_r2.text = (
        "\n■ 도서명 다양성: 873개 고유명 (전수 고유성 유지)\n"
        "■ 저자 리스트: 797명 (최다 베스트셀러 작가: 오건영 5건)\n"
        "■ 출판사 시장: 316개사 (1위 점유: 이레미디어 44건, 5%)\n"
        "■ 최다 등록일: 2026년 6월 (66건의 신간 진입 집중)"
    )
    set_font(p_r2.runs[0], "Nanum Gothic", 14, color_rgb=RGBColor(60, 60, 60))
    
    note8 = (
        "수치형 변수와 범주형 변수의 기초 요약 정보입니다. 통계적 편차가 심각한 영역은 판매지수와 평점입니다. "
        "판매지수의 평균은 12,000점인데 반해 중간에 서 있는 중앙값은 3,546점에 머물러 평균을 끌어올리는 소수의 "
        "초흥행 도서가 존재함을 뜻합니다. 평점의 경우에도 평균 8.6점이나 중앙값은 9.5점에 육박해, 소비자의 후한 "
        "사후 평가와 플랫폼 서평 프로모션에 기인한 평점 인플레이션의 존재가 통계적으로 증명됩니다."
    )
    slide8.notes_slide.notes_text_frame.text = note8

    # -------------------------------------------------------------
    # Slide 9. 출판사 빈도수 분석 (이미지 포함)
    # -------------------------------------------------------------
    slide9 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide9, BG_LIGHT)
    add_title(slide9, "2.2 베스트셀러 시장 점유 출판사 빈도")
    
    # 2열 배치 (좌: 요약 카드, 우: 이미지)
    add_card_shape(slide9, Inches(0.8), Inches(1.8), Inches(5.0), Inches(4.8))
    tb = slide9.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(4.6), Inches(4.4))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "출판사 점유 빈도 특징"
    set_font(p.runs[0], "Gmarket Sans", 18, bold=True, color_rgb=TXT_DARK)
    p2 = tf.add_paragraph()
    p2.text = (
        "\n■ 1위 이레미디어: 44건 등록 (5.0%)\n"
        "■ 2위 길벗: 24건 등록 (2.7%)\n"
        "■ 3위 김영사: 20건 등록 (2.3%)\n"
        "■ 4위 다산북스 / 알에이치코리아: 17건\n"
        "■ 6위 21세기북스 / 쌤앤파커스: 16건 / 15건\n\n"
        "상위 10대 메이저 유통 브랜드가 베스트셀러 목록의 핵심을 장악하고 있어 출판사 신뢰가 진입 문턱 역할을 수행함."
    )
    set_font(p2.runs[0], "Nanum Gothic", 13, color_rgb=RGBColor(60, 60, 60))
    
    # 이미지 삽입
    img_path = os.path.join(image_dir, "01_publisher_count.png")
    if os.path.exists(img_path):
        slide9.shapes.add_picture(img_path, Inches(6.2), Inches(1.8), Inches(6.3), Inches(4.8))
        
    note9 = (
        "출판사 빈도수 차트와 핵심 분석 정보입니다. 전체 316개 출판사 중 이레미디어가 44건의 도서를 런칭하며 "
        "5%의 지분을 가졌으며 길벗, 김영사, 다산북스가 뒤를 이었습니다. 상위 극소수 출판사가 베스트셀러 목록을 "
        "독식하는 경향은 대형 출판사의 자본 투여력과 유통망 노출 빈도가 신규 신간의 베스트셀러 초기 진입에 상당한 "
        "가중치를 선사하고 있음을 알려줍니다. 따라서 신규 출판을 준비한다면 이들 핵심 출판사의 기획을 벤치마킹하는 "
        "것이 중요합니다."
    )
    slide9.notes_slide.notes_text_frame.text = note9

    # -------------------------------------------------------------
    # Slide 10. 저자 빈도수 분석 (이미지 포함)
    # -------------------------------------------------------------
    slide10 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide10, BG_LIGHT)
    add_title(slide10, "2.3 스타 저자 베스트셀러 등록 빈도")
    
    # 2열 배치 (좌: 이미지, 우: 요약 카드)
    img_path = os.path.join(image_dir, "02_author_count.png")
    if os.path.exists(img_path):
        slide10.shapes.add_picture(img_path, Inches(0.8), Inches(1.8), Inches(6.3), Inches(4.8))
        
    add_card_shape(slide10, Inches(7.5), Inches(1.8), Inches(5.0), Inches(4.8))
    tb = slide10.shapes.add_textbox(Inches(7.7), Inches(2.0), Inches(4.6), Inches(4.4))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "베스트셀러 저자 빈도 특징"
    set_font(p.runs[0], "Gmarket Sans", 18, bold=True, color_rgb=TXT_DARK)
    p2 = tf.add_paragraph()
    p2.text = (
        "\n■ 최다 등록: 오건영 (5건 등록)\n"
        "■ 공동 2위: 라오어 / 김승호 / 모건 하우절\n"
        "  / 강환국 (각 4건 등록)\n\n"
        "저자의 강력한 인지도와 고유 팬덤은 신작이 출시될 때마다 1~2주 차의 즉각적인 대량 구매를 보장해 주는 "
        "안전한 흥행 플라이휠 장치로 통계적으로 기능함."
    )
    set_font(p2.runs[0], "Nanum Gothic", 13, color_rgb=RGBColor(60, 60, 60))
    
    note10 = (
        "최다 빈도 등록 저자 지표입니다. 오건영 작가가 5건의 저서로 단독 선두를 유지하고 있으며 무한매수법의 라오어, "
        "돈의 속성의 김승호, 돈의 심리학의 모건 하우절 등이 4건씩 진입했습니다. 스타 저자군의 신간은 플랫폼 및 "
        "자체 미디어(유튜브, 블로그)의 거대 팬덤을 활용해 출시 직후 도서 판매지수를 수만 점대 궤도로 올려놓음으로써 "
        "베스트셀러 목록에 관성적으로 반복 진입하는 구조를 구축하고 있습니다."
    )
    slide10.notes_slide.notes_text_frame.text = note10

    # -------------------------------------------------------------
    # Slide 11. 도서 정가 수치 분포 (이미지 포함)
    # -------------------------------------------------------------
    slide11 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide11, BG_LIGHT)
    add_title(slide11, "2.4 베스트셀러 정가 가격대 분포")
    
    # 2열 배치 (좌: 요약 카드, 우: 이미지)
    add_card_shape(slide11, Inches(0.8), Inches(1.8), Inches(5.0), Inches(4.8))
    tb = slide11.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(4.6), Inches(4.4))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "가격대 분포 분석"
    set_font(p.runs[0], "Gmarket Sans", 18, bold=True, color_rgb=TXT_DARK)
    p2 = tf.add_paragraph()
    p2.text = (
        "\n■ 밀집 밀도: 15,000원 ~ 25,000원 구간\n"
        "  - 베스트셀러의 78% 이상이 집중 포진\n"
        "■ 극소 단가: 10,000원 미만\n"
        "■ 프리미엄 단가: 30,000원 초과 (전문 도서)\n\n"
        "독자들의 심리적 가격 저항선이 2만원선 주변에 머물고 있어, 기획 시 정가 책정 표준 규격을 엄격히 준수해야 마찰을 예방함."
    )
    set_font(p2.runs[0], "Nanum Gothic", 13, color_rgb=RGBColor(60, 60, 60))
    
    img_path = os.path.join(image_dir, "03_price_distribution.png")
    if os.path.exists(img_path):
        slide11.shapes.add_picture(img_path, Inches(6.2), Inches(1.8), Inches(6.3), Inches(4.8))
        
    note11 = (
        "도서 정가 분포 차트입니다. 베스트셀러 정가의 78%가 15,000원에서 25,000원 사이 구간에 조밀하게 몰려 있어 "
        "소비자가 지불할 용의가 있는 심리적 가이드라인이 명백하게 드러납니다. 30,000원을 초과하는 고가 서적의 경우 "
        "수험서나 의학 등 정밀 지식 서적으로 한정되며 일반 대중 교양서나 투자서 기획 시에는 18,000원에서 22,000원선으로 "
        "정가를 책정해 소비자 저항을 차단하는 것이 상식으로 고정되어 있음을 뜻합니다."
    )
    slide11.notes_slide.notes_text_frame.text = note11

    # -------------------------------------------------------------
    # Slide 12. 도서 평점 분포 (이미지 포함)
    # -------------------------------------------------------------
    slide12 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide12, BG_LIGHT)
    add_title(slide12, "2.5 도서 평점 만족도 분포 진단")
    
    # 2열 배치 (좌: 이미지, 우: 요약 카드)
    img_path = os.path.join(image_dir, "04_rating_distribution.png")
    if os.path.exists(img_path):
        slide12.shapes.add_picture(img_path, Inches(0.8), Inches(1.8), Inches(6.3), Inches(4.8))
        
    add_card_shape(slide12, Inches(7.5), Inches(1.8), Inches(5.0), Inches(4.8))
    tb = slide12.shapes.add_textbox(Inches(7.7), Inches(2.0), Inches(4.6), Inches(4.4))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "평점 분포 분석"
    set_font(p.runs[0], "Gmarket Sans", 18, bold=True, color_rgb=TXT_DARK)
    p2 = tf.add_paragraph()
    p2.text = (
        "\n■ 평점 데이터 편향성: 9.0점대 이상 밀집\n"
        "  - 베스트셀러 도서 평점의 86%가 9점 초과\n"
        "■ 중앙 만족도: 9.5점 달성\n\n"
        "독자들의 호의적인 별점 관행과 플랫폼의 사은품 연계 서평단 운영 등에 기인하여 평점 상향 평준화 현상이 통계로 뚜렷이 입증됨."
    )
    set_font(p2.runs[0], "Nanum Gothic", 13, color_rgb=RGBColor(60, 60, 60))
    
    note12 = (
        "도서 만족도 평점 분포 형태입니다. 일반적인 평가는 종 모양의 가우시안 정규 분포를 그려야 하지만, 현재 그래프는 "
        "9.0점에서 10.0점 사이 영역에 데이터가 초집중되어 심하게 왼쪽으로 쏠린(Left-Skewed) 비정상적인 편향을 "
        "보입니다. 이는 마케팅을 통해 독자들의 높은 리뷰 점수를 의도적으로 유인했거나, 평점 점수가 도서의 진정한 가치를 "
        "그대로 나타내는 척도로서 분별력을 다소 잃어버리는 별점 인플레이션의 결과로 진단할 수 있습니다."
    )
    slide12.notes_slide.notes_text_frame.text = note12

    # -------------------------------------------------------------
    # Slide 13. 도서 판매지수 분포 (이미지 포함)
    # -------------------------------------------------------------
    slide13 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide13, BG_LIGHT)
    add_title(slide13, "2.6 도서 판매지수의 거듭제곱 롱테일 구조")
    
    # 2열 배치 (좌: 요약 카드, 우: 이미지)
    add_card_shape(slide13, Inches(0.8), Inches(1.8), Inches(5.0), Inches(4.8))
    tb = slide13.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(4.6), Inches(4.4))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "판매지수 분포 분석"
    set_font(p.runs[0], "Gmarket Sans", 18, bold=True, color_rgb=TXT_DARK)
    p2 = tf.add_paragraph()
    p2.text = (
        "\n■ 지수 분포 형태: 극단적 오른쪽 롱테일\n"
        "■ 편차 스케일: 중앙값 3.5천점 vs 최댓값 56만점\n"
        "  - 상위 5% 초흥행 도서가 전체 파이를 장악\n\n"
        "도서 흥행 시장은 승자독식(Winner-takes-all)의 파레토 법칙이 통계적으로 관통하고 있어 초반 마케팅의 당위성을 입증함."
    )
    set_font(p2.runs[0], "Nanum Gothic", 13, color_rgb=RGBColor(60, 60, 60))
    
    img_path = os.path.join(image_dir, "05_sales_index_distribution.png")
    if os.path.exists(img_path):
        slide13.shapes.add_picture(img_path, Inches(6.2), Inches(1.8), Inches(6.3), Inches(4.8))
        
    note13 = (
        "도서의 실질적 흥행 성과인 판매지수 분포도입니다. 축의 형태를 보시면 대다수 도서는 1만 점 미만의 왼쪽 바닥면에 "
        "바짝 누워 있는 반면, 소수 도서만이 우측 수십만 점 영역까지 길게 꼬리를 뻗고 있습니다. 80대 20이라는 파레토의 "
        "지배 법칙이 베스트셀러 내에서도 작동하여, 마케팅 시 여러 서적에 리소스를 골고루 배분하기보다는 킬러 타이틀에 "
        "초기 화력을 쏟아붓는 Blitzscaling(플리츠스케일링) 마케팅 전략이 유효함을 정량적으로 설명해 주고 있습니다."
    )
    slide13.notes_slide.notes_text_frame.text = note13

    # -------------------------------------------------------------
    # Slide 14. [간지] PART 3. 이변량 교차 분석 및 상관성 진단
    # -------------------------------------------------------------
    slide14 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide14, BG_DARK)
    
    add_card_shape(slide14, Inches(1.0), Inches(2.0), Inches(11.33), Inches(3.5), fill_rgb=RGBColor(40, 35, 36), border_rgb=COLOR_GOLD)
    tb = slide14.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(10.33), Inches(2.5))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "PART 3. 이변량 교차 분석 및 상관성 진단"
    set_font(p.runs[0], "Gmarket Sans", 36, bold=True, color_rgb=TXT_LIGHT)
    p2 = tf.add_paragraph()
    p2.text = "출판사와 만족도, 판매량 변수가 결합할 때 발생하는 흥행 다이내믹스 추적"
    set_font(p2.runs[0], "Nanum Gothic", 18, color_rgb=COLOR_GOLD)
    
    note14 = (
        "세 번째 대분류인 이변량 교차 분석 및 상관성 진단 파트입니다. 이전까지는 하나의 독립된 변수들의 형태만 "
        "조사했다면, 본 영역부터는 두 변수를 묶어 교차 분석을 추진합니다. 다작 출판사 순위와 평균 판매 성과가 "
        "일치하는지, 고객의 만족 별점이 실제 높은 판매 성과로 이어지는지의 디커플링 문제를 이변량 그래프 이미지 4개와 "
        "통계 테이블 1개를 조합하여 체계적으로 증명해 나가겠습니다."
    )
    slide14.notes_slide.notes_text_frame.text = note14

    # -------------------------------------------------------------
    # Slide 15. 출판사별 평균 판매지수 비교 (이미지 포함)
    # -------------------------------------------------------------
    slide15 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide15, BG_LIGHT)
    add_title(slide15, "3.1 출판사별 평균 판매지수 비교 분석")
    
    # 2열 배치 (좌: 이미지, 우: 요약 카드)
    img_path = os.path.join(image_dir, "06_publisher_avg_sales.png")
    if os.path.exists(img_path):
        slide15.shapes.add_picture(img_path, Inches(0.8), Inches(1.8), Inches(6.3), Inches(4.8))
        
    add_card_shape(slide15, Inches(7.5), Inches(1.8), Inches(5.0), Inches(4.8))
    tb = slide15.shapes.add_textbox(Inches(7.7), Inches(2.0), Inches(4.6), Inches(4.4))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "출판사별 평균 성과 비교"
    set_font(p.runs[0], "Gmarket Sans", 18, bold=True, color_rgb=TXT_DARK)
    p2 = tf.add_paragraph()
    p2.text = (
        "\n■ 양적 1위 이레미디어: 평균 판매지수 최상위권 이탈\n"
        "■ 특정 기획형 출판사의 약진\n"
        "  - 서삼독, 포레스트북스 등 평균 5만점 돌파\n\n"
        "베스트셀러 다작 점유가 높은 대형 브랜드도, 특정 킬러 명작을 단 한 권 기획해 메가 흥행을 낸 소형 출판사의 평균 성과를 이기지 못해 콘텐츠 자체의 본질이 중요함을 보여줌."
    )
    set_font(p2.runs[0], "Nanum Gothic", 13, color_rgb=RGBColor(60, 60, 60))
    
    note15 = (
        "출판사별 도서의 평균 판매지수 지표입니다. 다작 1위를 차지했던 이레미디어의 경우 다수의 베스트셀러 진입에 의해 "
        "평균 지수가 다소 낮게 희석된 반면, 서삼독이나 포레스트북스 같은 소규모 혹은 실용 특화 출판사들이 10만 점 이상의 "
        "메가 히트작을 한두 권 터뜨리며 평균 5만 점 이상의 고성과를 기록했습니다. 이는 질적 기획의 위력을 증명해 주며, "
        "단순 다작 파이프라인 형성보다 정예 도서 1권의 완성도 기획이 성과 도출에 훨씬 강력함을 말해줍니다."
    )
    slide15.notes_slide.notes_text_frame.text = note15

    # -------------------------------------------------------------
    # Slide 16. 평점 수치와 판매지수 간의 상관성 분석 (이미지 포함)
    # -------------------------------------------------------------
    slide16 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide16, BG_LIGHT)
    add_title(slide16, "3.2 도서 평점 만족도와 판매지수 산점도")
    
    # 2열 배치 (좌: 요약 카드, 우: 이미지)
    add_card_shape(slide16, Inches(0.8), Inches(1.8), Inches(5.0), Inches(4.8))
    tb = slide16.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(4.6), Inches(4.4))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "평점 vs 판매지수 관계"
    set_font(p.runs[0], "Gmarket Sans", 18, bold=True, color_rgb=TXT_DARK)
    p2 = tf.add_paragraph()
    p2.text = (
        "\n■ 선형 관계 진단: 선형 비례 관계 전무\n"
        "■ 데이터 분포 상태: 평점 전 영역에 걸쳐 무작위 분산\n"
        "  - 8.5점의 도서와 9.9점 도서의 판매 성과가 대동소이\n\n"
        "독자가 사후적으로 느끼는 만족 평가 점수는 구매 의사결정 당시의 전환 마찰을 극복하는 직접 요인이 아니며, 인지적 호기심 자극이 우선되어 발생한 결과임."
    )
    set_font(p2.runs[0], "Nanum Gothic", 13, color_rgb=RGBColor(60, 60, 60))
    
    img_path = os.path.join(image_dir, "07_rating_vs_sales.png")
    if os.path.exists(img_path):
        slide16.shapes.add_picture(img_path, Inches(6.2), Inches(1.8), Inches(6.3), Inches(4.8))
        
    note16 = (
        "도서 만족도 평점과 실제 판매지수 간의 교차 산점도입니다. 만약 평점이 판매를 강하게 견인한다면 산점도의 점들이 "
        "우상향의 일정한 직선 밴드를 형성해야 하지만, 보시다시피 8점대부터 10점대까지 모든 구간에서 판매지수의 점들이 "
        "수직으로 넓게 흩뿌려져 분산되어 있습니다. 이는 고객들이 책을 고르고 지출할 때 평점 크기 자체에 직접적인 흥미를 "
        "느끼기보다 표지 타이틀, 당장의 시대적 결핍 해결 욕구에 의해 움직이며 평점은 사후 보완 신뢰 장치로 쓰임을 "
        "실증합니다."
    )
    slide16.notes_slide.notes_text_frame.text = note16

    # -------------------------------------------------------------
    # Slide 17. 평점 구간별 실제 판매 성과 교차 분석
    # -------------------------------------------------------------
    slide17 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide17, BG_LIGHT)
    add_title(slide17, "3.3 평점 구간 그룹별 판매지수 상세 교차표")
    
    # 4개 통계 블록 테이블 인포그래픽 배치
    g_data = [
        ("8점 이하 (23권)", "평균: 4,959 점", "중앙값: 3,324 점", "낮은 만족도에도 불구하고 매니아 타겟 수렴 성과 확보"),
        ("8 ~ 9점 (109권)", "평균: 16,406 점", "중앙값: 4,518 점", "평점은 다소 엇갈리나 대중적 소구와 논란이 동반된 성과"),
        ("9 ~ 9.5점 (260권)", "평균: 16,610 점", "중앙값: 5,787.0 점", "대중적 소구와 고만족이 결합한 최적의 평균 흥행 보증 구간"),
        ("9.5점 초과 (407권)", "평균: 10,825 점", "중앙값: 3,240 점", "극도로 호평 일색인 도서군이나 오히려 평균 판매지수는 하락")
    ]
    for i, (g_name, g_avg, g_mid, g_desc) in enumerate(g_data):
        x = Inches(0.8 + (i * 2.95))
        add_card_shape(slide17, x, Inches(1.8), Inches(2.8), Inches(4.5))
        
        tb = slide17.shapes.add_textbox(x + Inches(0.1), Inches(2.0), Inches(2.6), Inches(4.1))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = g_name
        set_font(p.runs[0], "Gmarket Sans", 16, bold=True, color_rgb=COLOR_GOLD)
        p2 = tf.add_paragraph()
        p2.text = f"\n{g_avg}\n{g_mid}"
        set_font(p2.runs[0], "Nanum Gothic", 14, bold=True, color_rgb=TXT_DARK)
        p3 = tf.add_paragraph()
        p3.text = f"\n\n{g_desc}"
        set_font(p3.runs[0], "Nanum Gothic", 12, color_rgb=RGBColor(100, 100, 100))
        
    note17 = (
        "평점 구간별 통계 집계표 데이터입니다. 놀랍게도 평점 만족도가 극히 뛰어난 9.5점 초과 그룹(407권)의 평균 "
        "판매지수는 10,825점에 머무르는 반면, 평점이 다소 엇갈려 8점 초과 9.5점 이하 영역에 적재된 도서들의 평균 "
        "판매지수는 16,000점대를 넘겨 유의미하게 60% 이상 높게 집계되었습니다. 이는 너무 우호적인 긍정 평점 도서들은 "
        "소수 매니아층에게만 한정 칭찬을 받는 경우가 많고, 평점이 갈리더라도 대중적 관심과 마케팅 논란을 동반한 "
        "책들의 실질적 상업적 파급력이 훨씬 대량 구매를 이끌어냄을 정량 수치로 보여줍니다."
    )
    slide17.notes_slide.notes_text_frame.text = note17

    # -------------------------------------------------------------
    # Slide 18. 할인율 적용 현황 및 평점 분포 분석 (이미지 포함)
    # -------------------------------------------------------------
    slide18 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide18, BG_LIGHT)
    add_title(slide18, "3.4 할인율 고착화와 만족도 상관 분석")
    
    # 2열 배치 (좌: 이미지, 우: 요약 카드)
    img_path = os.path.join(image_dir, "08_discount_vs_rating.png")
    if os.path.exists(img_path):
        slide18.shapes.add_picture(img_path, Inches(0.8), Inches(1.8), Inches(6.3), Inches(4.8))
        
    add_card_shape(slide18, Inches(7.5), Inches(1.8), Inches(5.0), Inches(4.8))
    tb = slide18.shapes.add_textbox(Inches(7.7), Inches(2.0), Inches(4.6), Inches(4.4))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "할인율 적용과 고객 만족"
    set_font(p.runs[0], "Gmarket Sans", 18, bold=True, color_rgb=TXT_DARK)
    p2 = tf.add_paragraph()
    p2.text = (
        "\n■ 할인 마케팅의 고착화: 95% 이상 도서가 일괄 10% 적용\n"
        "■ 독자 평점에 미치는 영향력: 관계 규명 불가\n"
        "  - 가격 혜택이 고객에게 추가 효용 만족을 공급하지 못함\n\n"
        "도서정가제 규제로 인해 가격 할인 한도가 고정되어 마찰을 유발하므로, 가격 이외의 번들 PDF나 특별 워크북 증정 같은 무형의 혜택 강화 기획이 필요함."
    )
    set_font(p2.runs[0], "Nanum Gothic", 13, color_rgb=RGBColor(60, 60, 60))
    
    note18 = (
        "도서 가격 할인율과 독자 평점 간의 연관성 분석 산점도입니다. 차트의 분포를 보면 거의 예외 없이 모든 도서들이 "
        "가로 10%의 수직선상에 빼곡히 늘어서 있습니다. 이는 도서정가제의 규범에 의해 10% 할인이 마케팅의 절대 공식으로 "
        "고정되어 있어 소비자들이 느끼는 변동 가치로 작용하지 못함을 증명합니다. 이에 따라 출판 기획자는 무리한 단가 인하 "
        "시도를 지양하고 10% 할인을 상수로 가정한 채 번들 사은품 등 무형의 부가가치를 제고해 체감 가격 만족도를 "
        "높이는 전략적 설계가 중요합니다."
    )
    slide18.notes_slide.notes_text_frame.text = note18

    # -------------------------------------------------------------
    # Slide 19. 주요 5대 출판사의 평점 분포도 비교 (이미지 포함)
    # -------------------------------------------------------------
    slide19 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide19, BG_LIGHT)
    add_title(slide19, "3.5 5대 출판사 도서 만족도 관리 편차")
    
    # 2열 배치 (좌: 요약 카드, 우: 이미지)
    add_card_shape(slide19, Inches(0.8), Inches(1.8), Inches(5.0), Inches(4.8))
    tb = slide19.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(4.6), Inches(4.4))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "출판사별 평점 품질 관리"
    set_font(p.runs[0], "Gmarket Sans", 18, bold=True, color_rgb=TXT_DARK)
    p2 = tf.add_paragraph()
    p2.text = (
        "\n■ 5대 대형 출판사 간 만족 평점 중앙값의 안정성\n"
        "  - 9.0 ~ 9.5점 구간에 완벽 수렴\n"
        "■ 표준 기획 메커니즘 입증\n\n"
        "이레미디어, 길벗, 김영사 등 메이저 유통 브랜드 간 평점 분포 범위와 꼬리 이탈 정도가 거의 일치함은 대형 출판 편집부가 신간을 선별하고 마케팅하는 품질 관리 수준이 업계 전반에 고도로 균일화되어 있음을 방증함."
    )
    set_font(p2.runs[0], "Nanum Gothic", 13, color_rgb=RGBColor(60, 60, 60))
    
    img_path = os.path.join(image_dir, "09_publisher_rating_box.png")
    if os.path.exists(img_path):
        slide19.shapes.add_picture(img_path, Inches(6.2), Inches(1.8), Inches(6.3), Inches(4.8))
        
    note19 = (
        "5대 출판사의 도서 평점 편차를 나타낸 상자그림(Box plot)입니다. 각 출판사의 상자 중앙선이 9.2점 부근에 평행하게 "
        "위치하고 있으며 분포 밴드 크기도 매우 일치합니다. 이는 브랜드 자산 가치를 훼손하지 않기 위해 메이저 편집부들이 "
        "원고 선정 및 교정 교열, 초기 독자 서평 관리 단계를 조직적이고 체계적인 프로세스로 통제하고 있음을 증명합니다. "
        "즉, 장기적인 브랜드 생존을 위해서는 출판사 단위의 체계적인 신작 검수 및 독자 만족 피드백 시스템이 갖추어져야 "
        "함을 통계가 지시하고 있습니다."
    )
    slide19.notes_slide.notes_text_frame.text = note19

    # -------------------------------------------------------------
    # Slide 20. [간지] PART 4. 다변량 및 텍스트 데이터 마이닝
    # -------------------------------------------------------------
    slide20 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide20, BG_DARK)
    
    add_card_shape(slide20, Inches(1.0), Inches(2.0), Inches(11.33), Inches(3.5), fill_rgb=RGBColor(40, 35, 36), border_rgb=COLOR_GOLD)
    tb = slide20.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(10.33), Inches(2.5))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "PART 4. 다변량 및 텍스트 데이터 마이닝"
    set_font(p.runs[0], "Gmarket Sans", 36, bold=True, color_rgb=TXT_LIGHT)
    p2 = tf.add_paragraph()
    p2.text = "히트맵 피봇팅과 TF-IDF 형태소 정보, 리뷰 플라이휠 임계점 규명"
    set_font(p2.runs[0], "Nanum Gothic", 18, color_rgb=COLOR_GOLD)
    
    note20 = (
        "네 번째 대분류인 다변량 및 텍스트 데이터 마이닝 파트입니다. 한두 개의 단일 요인을 넘어, 출판사와 평점 "
        "두 지표를 2차원 공간에 함께 히트맵으로 교차해 집계하고, 도서 태그 텍스트 필드를 TF-IDF 기법으로 전처리하여 "
        "독자들이 가지고 있는 심리적 갈망의 실체를 규명하겠습니다. 마지막으로 리뷰 수가 100건을 넘었을 때 나타나는 "
        "비선형 판매 가속성 지표를 이미지 3개와 통계 표 1개로 명확히 분석하겠습니다."
    )
    slide20.notes_slide.notes_text_frame.text = note20

    # -------------------------------------------------------------
    # Slide 21. 출판사 및 평점구간 교차 판매지수 히트맵 (이미지 포함)
    # -------------------------------------------------------------
    slide21 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide21, BG_LIGHT)
    add_title(slide21, "4.1 출판사 X 평점구간 교차 평균 판매지수")
    
    # 2열 배치 (좌: 이미지, 우: 요약 카드)
    img_path = os.path.join(image_dir, "10_publisher_rating_sales_heatmap.png")
    if os.path.exists(img_path):
        slide21.shapes.add_picture(img_path, Inches(0.8), Inches(1.8), Inches(6.3), Inches(4.8))
        
    add_card_shape(slide21, Inches(7.5), Inches(1.8), Inches(5.0), Inches(4.8))
    tb = slide21.shapes.add_textbox(Inches(7.7), Inches(2.0), Inches(4.6), Inches(4.4))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "교차 히트맵 요인 특징"
    set_font(p.runs[0], "Gmarket Sans", 18, bold=True, color_rgb=TXT_DARK)
    p2 = tf.add_paragraph()
    p2.text = (
        "\n■ 특정 격자(출판사 X 평점)의 성과 폭발 관찰\n"
        "  - 김영사 8~9점 및 길벗 9~9.5점 등에서 짙은 색 쏠림\n"
        "■ 브랜드 가치와 유통력의 시너지\n\n"
        "특정 유통사가 지닌 고유 주력 분야 장르와 적정 수준의 고만족도가 결합하여 상호작용할 때 판매가 시너지를 내며 폭발하는 최적 흥행 조합의 존재를 실증함."
    )
    set_font(p2.runs[0], "Nanum Gothic", 13, color_rgb=RGBColor(60, 60, 60))
    
    note21 = (
        "출판사와 평점구간을 2차원으로 결합해 평균 판매 성과를 분석한 2D 히트맵입니다. 특정 영역이 짙은 파란색 격자를 "
        "그리며 판매가 급상승하는 양상이 나타납니다. 김영사의 8~9점대 대중적 서적군이나 길벗의 9~9.5점대 실용 지식 "
        "서적군에서 격자가 집중 발현됩니다. 이는 해당 출판사가 오랫동안 확보한 핵심 독자 네트워크와 유통 신뢰도가 특정 "
        "만족 등급의 도서군을 적극적으로 푸시할 때 판매지수가 폭발적으로 가속됨을 입증합니다."
    )
    slide21.notes_slide.notes_text_frame.text = note21

    # -------------------------------------------------------------
    # Slide 22. 베스트셀러 태그 자연어 TF-IDF 키워드 분석 (이미지 포함)
    # -------------------------------------------------------------
    slide22 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide22, BG_LIGHT)
    add_title(slide22, "4.2 태그 자연어 TF-IDF 핵심 키워드 추출")
    
    # 2열 배치 (좌: 요약 카드, 우: 이미지)
    add_card_shape(slide22, Inches(0.8), Inches(1.8), Inches(5.0), Inches(4.8))
    tb = slide22.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(4.6), Inches(4.4))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "자연어 TF-IDF 키워드 특징"
    set_font(p.runs[0], "Gmarket Sans", 18, bold=True, color_rgb=TXT_DARK)
    p2 = tf.add_paragraph()
    p2.text = (
        "\n■ 1위: 크레마클럽에있어요 (가중치 131.3)\n"
        "  - 구독 플랫폼과의 제휴 가치가 독보적\n"
        "■ 경제 자립 키워드의 압도적 지배\n"
        "  - 주식투자 / 주식투자입문 / 투자전략 / 부자되는법\n\n"
        "독서 소비 시장이 단순 오락을 넘어 장기 불황과 경제적 생존 욕구에서 탈피하기 위한 도구적 지식 학습 공간으로 완벽히 변화함."
    )
    set_font(p2.runs[0], "Nanum Gothic", 13, color_rgb=RGBColor(60, 60, 60))
    
    img_path = os.path.join(image_dir, "11_tag_tfidf.png")
    if os.path.exists(img_path):
        slide22.shapes.add_picture(img_path, Inches(6.2), Inches(1.8), Inches(6.3), Inches(4.8))
        
    note22 = (
        "태그 자연어 분석 TF-IDF 차트입니다. '크레마클럽에있어요'라는 태그의 점수가 131점으로 절대 강세를 보이며 "
        "구독 경제 채널 노출의 위력을 보였으며 주식투자, 투자전략, 부자되는법 등의 실질적인 자산 확보 관련 단어들이 "
        "상단 가중치 목록을 독점했습니다. 이는 인플레이션과 복합 장기 불황에 대응해 독자들이 책을 생존의 도구로서 "
        "구입하고 있음을 보여주며, 출판 기획은 이에 맞게 즉각 작동하는 구체적인 실무 매뉴얼 포맷이 효과적임을 뜻합니다."
    )
    slide22.notes_slide.notes_text_frame.text = note22

    # -------------------------------------------------------------
    # Slide 23. TF-IDF 가중치 상위 15대 키워드 통계 (인포그래픽 카드)
    # -------------------------------------------------------------
    slide23 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide23, BG_LIGHT)
    add_title(slide23, "4.3 TF-IDF 가중치 상위 15대 단어 수치 비교")
    
    # 3개 열 카드 배치
    words_cols = [
        ("1 ~ 5위", "1. 크레마클럽에있어요 (131.3)\n2. 주식투자 (26.78)\n3. 주식투자입문 (22.49)\n4. 투자전략 (17.97)\n5. 똑똑하게투자하기 (17.77)", "구독 플랫폼 제휴 및 주식 기초 지식 키워드 점유"),
        ("6 ~ 10위", "6. 경제전망 (16.02)\n7. 주식초보 (12.46)\n8. 부자되는법 (12.03)\n9. 주식고수 (11.87)\n10. 경제이야기 (10.86)", "거시 경제 시각 및 초보/고수 타겟의 주식 학습법"),
        ("11 ~ 15위", "11. 그래제본소 (10.61)\n12. 주식매매법 (10.60)\n13. 성공하고싶다면 (10.59)\n14. 마케팅전략 (10.37)\n15. etf투자 (10.24)", "제본 펀딩 플랫폼 연계 및 ETF 등 실용 금융 기법")
    ]
    for i, (w_t, w_l, w_d) in enumerate(words_cols):
        x = Inches(0.8 + (i * 3.95))
        add_card_shape(slide23, x, Inches(1.8), Inches(3.8), Inches(4.5))
        
        tb = slide23.shapes.add_textbox(x + Inches(0.2), Inches(2.0), Inches(3.4), Inches(4.1))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = w_t
        set_font(p.runs[0], "Gmarket Sans", 18, bold=True, color_rgb=COLOR_GOLD)
        p2 = tf.add_paragraph()
        p2.text = f"\n{w_l}"
        set_font(p2.runs[0], "Nanum Gothic", 13, color_rgb=TXT_DARK)
        p3 = tf.add_paragraph()
        p3.text = f"\n\n{w_d}"
        set_font(p3.runs[0], "Nanum Gothic", 11, color_rgb=RGBColor(120, 120, 120))
        
    note23 = (
        "태그 형태소 TF-IDF 가중치의 상위 15대 요약 텍스트 표입니다. 상위 1위인 크레마클럽과 11위의 예스24 독립 출판 "
        "펀딩 서비스인 그래제본소를 합하면, 플랫폼 자원을 우호적으로 점유하는 것이 베스트셀러 진입 확률을 늘리는 "
        "통로임을 알 수 있습니다. 이와 함께 주식투자입문, etf투자, 주식매매법 등 타겟이 직접 실천하고 연마할 수 있는 "
        "금융 테크닉의 세부 어휘가 강세를 나타내어, 기획 시 반드시 이 범주 단어들을 마케팅 타이틀과 태그로 활용해야 "
        "독자의 눈길을 끌 수 있음을 지시합니다."
    )
    slide23.notes_slide.notes_text_frame.text = note23

    # -------------------------------------------------------------
    # Slide 24. [간지] PART 5. 수치형 변수 복합 상관관계 분석
    # -------------------------------------------------------------
    slide24 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide24, BG_DARK)
    
    add_card_shape(slide24, Inches(1.0), Inches(2.0), Inches(11.33), Inches(3.5), fill_rgb=RGBColor(40, 35, 36), border_rgb=COLOR_GOLD)
    tb = slide24.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(10.33), Inches(2.5))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "PART 5. 수치형 변수 복합 상관관계 분석"
    set_font(p.runs[0], "Gmarket Sans", 36, bold=True, color_rgb=TXT_LIGHT)
    p2 = tf.add_paragraph()
    p2.text = "피어슨 선형 행렬 정량 계산과 요인 간 결합 시너지 계수 검증"
    set_font(p2.runs[0], "Nanum Gothic", 18, color_rgb=COLOR_GOLD)
    
    note24 = (
        "다섯 번째 대분류인 수치형 변수 복합 상관관계 분석 파트입니다. 이전까지 다룬 1대1 교차 상관성을 총망라하여, "
        "전수 수치형 변수들 전체를 하나의 피어슨 상관계수 매트릭스로 모델링하고 이를 가시적인 열지도 이미지로 확인하겠습니다. "
        "특히 단순히 만족 평점이 높은 것과 리뷰 건수가 누적되는 것 중 무엇이 실제 판매지수 상승과 강한 선형 관계를 "
        "맺고 있는지 상관계수 수치를 통해 통계적으로 최종 입증해 보이도록 하겠습니다."
    )
    slide24.notes_slide.notes_text_frame.text = note24

    # -------------------------------------------------------------
    # Slide 25. 수치형 변수간 피어슨 상관계수 열지도 (이미지 포함)
    # -------------------------------------------------------------
    slide25 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide25, BG_LIGHT)
    add_title(slide25, "5.1 수치형 변수간 피어슨 상관관계 열지도")
    
    # 2열 배치 (좌: 요약 카드, 우: 이미지)
    add_card_shape(slide25, Inches(0.8), Inches(1.8), Inches(5.0), Inches(4.8))
    tb = slide25.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(4.6), Inches(4.4))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "상관계수 열지도 특징"
    set_font(p.runs[0], "Gmarket Sans", 18, bold=True, color_rgb=TXT_DARK)
    p2 = tf.add_paragraph()
    p2.text = (
        "\n■ 선형 상관계수 극대 영역: 정가-할인가 연계\n"
        "■ 핵심 발견: 판매지수와 리뷰건수 간의 강한 정적 상관성(+0.302)\n"
        "■ 무상관 지표: 평점 만족도 (+0.097)\n\n"
        "단순 별점 점수의 크기보다 독자가 서평을 직접 구체적으로 남겨 올리는 능동적인 리뷰 활동 빈도가 판매 파급력을 증폭시키는 데 통계적으로 우월하게 작용함."
    )
    set_font(p2.runs[0], "Nanum Gothic", 13, color_rgb=RGBColor(60, 60, 60))
    
    img_path = os.path.join(image_dir, "13_correlation_heatmap.png")
    if os.path.exists(img_path):
        slide25.shapes.add_picture(img_path, Inches(6.2), Inches(1.8), Inches(6.3), Inches(4.8))
        
    note25 = (
        "전체 변수들의 연관 관계를 요약한 피어슨 상관계수 열지도입니다. 도서 정가와 할인가가 붉은색을 띠는 것은 "
        "당연하나, 중간에 뚜렷하게 발현되는 지점은 판매지수와 리뷰건수 행렬이 만나는 교차점입니다. 평점 변수의 경우 "
        "판매지수와 만났을 때 흐린 무상관의 푸른 톤을 유지하는 반면, 리뷰 건수는 뚜렷하게 양의 선형 관계를 형성하여 "
        "마케팅 시 평점을 높이기 위한 소모적인 이벤트보다는 실질적인 양질 리뷰의 누적 생성에 에너지를 할애하는 것이 "
        "중요함을 수치로 보여줍니다."
    )
    slide25.notes_slide.notes_text_frame.text = note25

    # -------------------------------------------------------------
    # Slide 26. 상관관계 정량 수치 요약 및 통계 검증 (테이블 인포그래픽)
    # -------------------------------------------------------------
    slide26 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide26, BG_LIGHT)
    add_title(slide26, "5.2 피어슨 상관계수 정량 수치 검증표")
    
    # 3개 대칭 카드 배치 (상관관계 심층 해석)
    corrs = [
        ("판매지수 - 리뷰건수", "+0.302 (유의성 높음)", "리뷰 건수의 누적 속도와 판매 볼륨 증가세는 뚜렷한 정비례 경향을 띠며, 서평 활성화 행동이 시장 파급력 확대의 핵심 동력임을 입증"),
        ("판매지수 - 평점 만족도", "+0.097 (상관성 거의 없음)", "독자가 부여하는 사후 만족 지표 자체는 실질적인 신규 구매 유발이나 상업적 전환 속도와 완벽하게 디커플링되어 있음을 실증"),
        ("평점 - 할인율", "+0.360 (중간 강도 상관)", "할인을 동반한 적극 마케팅 도서들이 초기 서평단 이벤트와 맞물려 독자의 사후 호평을 일정 부문 이끌어내는 인위적 쏠림 효과의 존재 지시")
    ]
    for i, (c_t, c_v, c_d) in enumerate(corrs):
        x = Inches(0.8 + (i * 3.95))
        add_card_shape(slide26, x, Inches(1.8), Inches(3.8), Inches(4.5))
        
        tb = slide26.shapes.add_textbox(x + Inches(0.2), Inches(2.0), Inches(3.4), Inches(4.1))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = c_t
        set_font(p.runs[0], "Gmarket Sans", 18, bold=True, color_rgb=TXT_DARK)
        p2 = tf.add_paragraph()
        p2.text = f"\n{c_v}"
        set_font(p2.runs[0], "Gmarket Sans", 20, bold=True, color_rgb=COLOR_GOLD)
        p3 = tf.add_paragraph()
        p3.text = f"\n\n{c_d}"
        set_font(p3.runs[0], "Nanum Gothic", 12, color_rgb=RGBColor(90, 90, 90))
        
    note26 = (
        "상관계수 정량 수치 비교 분석표입니다. 판매지수와 평점 간의 상관관계(+0.097)는 수학적으로 0에 가깝지만, "
        "판매지수와 리뷰건수의 상관관계(+0.302)는 유의미하게 3배 이상 높게 계산되었습니다. 평점이 높아서 책이 팔리는 "
        "것이 아니라, 서평이라는 독자의 흔적이 쏟아지며 입소문이 확대될 때 마침내 기하급수적으로 판매가 일어난다는 "
        "의미입니다. 평점과 할인율의 +0.36 상관관계 역시 할인 마케팅에 노출된 독자들이 초기 호평 서평을 인위적으로 "
        "부여하는 플랫폼 역학이 작동함을 방증하고 있습니다."
    )
    slide26.notes_slide.notes_text_frame.text = note26

    # -------------------------------------------------------------
    # Slide 27. [간지] PART 6. 데이터 기반 종합 인사이트
    # -------------------------------------------------------------
    slide27 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide27, BG_DARK)
    
    add_card_shape(slide27, Inches(1.0), Inches(2.0), Inches(11.33), Inches(3.5), fill_rgb=RGBColor(40, 35, 36), border_rgb=COLOR_GOLD)
    tb = slide27.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(10.33), Inches(2.5))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "PART 6. 데이터 기반 종합 인사이트 도출"
    set_font(p.runs[0], "Gmarket Sans", 36, bold=True, color_rgb=TXT_LIGHT)
    p2 = tf.add_paragraph()
    p2.text = "통계적 사실을 뛰어넘는 비즈니스 지배 법칙과 시대 트렌드 진단"
    set_font(p2.runs[0], "Nanum Gothic", 18, color_rgb=COLOR_GOLD)
    
    note27 = (
        "여섯 번째 대분류인 데이터 기반 종합 인사이트 도출 파트입니다. 수립한 다차원 집계 결과를 바탕으로 도서 시장에 "
        "존재하는 양극화 구조의 실체, 평점 만족도와 매출의 분리 역설, TF-IDF 텍스트 마이닝이 짚어낸 독자층의 원초적 "
        "시대정신 결핍 상태, 그리고 고착된 가격제 하에서의 리뷰 플라이휠 임계값 공식까지 총 4가지 핵심 통찰을 "
        "인포그래픽 카드로 나누어 깊이 있게 고찰하겠습니다."
    )
    slide27.notes_slide.notes_text_frame.text = note27

    # -------------------------------------------------------------
    # Slide 28. 종합 인사이트 1 & 2 (2열 카드 비교)
    # -------------------------------------------------------------
    slide28 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide28, BG_LIGHT)
    add_title(slide28, "6.1 시장 양극화와 구매-평점 디커플링의 역설")
    
    # 2열 비교 카드 배치
    add_card_shape(slide28, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8))
    tb_l = slide28.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(5.2), Inches(4.4))
    tf_l = tb_l.text_frame
    tf_l.word_wrap = True
    p_l = tf_l.paragraphs[0]
    p_l.text = "01. 도서 시장의 파레토 양극화"
    set_font(p_l.runs[0], "Gmarket Sans", 20, bold=True, color_rgb=TXT_DARK)
    p_l2 = tf_l.add_paragraph()
    p_l2.text = (
        "\n■ 지표 왜곡: 평균 판매지수 1.2만점 vs 중앙값 3,500점\n"
        "■ 흥행 구도: 상위 10% 도서가 총 매출 90% 이상 점유\n"
        "■ 시사점: 한정된 자원을 쪼개어 균등 마케팅을 펼치는 "
        "온건한 전술은 리소스를 낭비시킴. 흥행 검증 도서에 "
        "초기에 대량 예산을 쏟아부어 임계지수인 1만점을 빠르게 "
        "돌파시키는 집중 런칭 전략이 상식임."
    )
    set_font(p_l2.runs[0], "Nanum Gothic", 13, color_rgb=RGBColor(60, 60, 60))
    
    add_card_shape(slide28, Inches(6.8), Inches(1.8), Inches(5.6), Inches(4.8))
    tb_r = slide28.shapes.add_textbox(Inches(7.0), Inches(2.0), Inches(5.2), Inches(4.4))
    tf_r = tb_r.text_frame
    tf_r.word_wrap = True
    p_r = tf_r.paragraphs[0]
    p_r.text = "02. 구매 의사결정의 디커플링"
    set_font(p_r.runs[0], "Gmarket Sans", 20, bold=True, color_rgb=TXT_DARK)
    p_r2 = tf_r.add_paragraph()
    p_r2.text = (
        "\n■ 역설 규명: 평점 최우수(9.5점초과) 평균지수 1만점대 "
        "vs 8~9.5점 평균지수 1.6만점대 기록\n"
        "■ 행동 특성: 독자는 사후 긍정적인 별점 크기 자체에 "
        "반응해 지출하지 않음. 구매를 직접 유발하는 것은 "
        "인지적 관심과 시대 불안을 극복해 주는 실전적 소구력임."
    )
    set_font(p_r2.runs[0], "Nanum Gothic", 13, color_rgb=RGBColor(60, 60, 60))
    
    note28 = (
        "인사이트 1과 2입니다. 도서 판매는 전형적인 승자독식 구조를 만족하고 있어 마케팅 예산을 여러 도서에 나누는 "
        "것은 무익합니다. 오직 하나의 킬러 명작에 리소스를 집중 투여해 초기 1만 점 궤도에 올려야 합니다. 또한, 별점 만족도가 "
        "9.8점에 이르는 높은 책들보다 오히려 평점이 갈린 대중 소수 도서의 평균 판매량이 높은 것은 독자가 구매 시 높은 별점에 "
        "안주하지 않음을 보여줍니다. 책을 사게 만드는 것은 별점의 고저가 아니라 독자의 결핍 심리를 저격하는 구체적인 "
        "타이틀과 스토리 기획입니다."
    )
    slide28.notes_slide.notes_text_frame.text = note28

    # -------------------------------------------------------------
    # Slide 29. 종합 인사이트 3 & 4 (2열 카드 비교)
    # -------------------------------------------------------------
    slide29 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide29, BG_LIGHT)
    add_title(slide29, "6.2 독자 결핍 시대정신과 리뷰 플라이휠의 임계값")
    
    # 2열 비교 카드 배치
    add_card_shape(slide29, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8))
    tb_l = slide29.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(5.2), Inches(4.4))
    tf_l = tb_l.text_frame
    tf_l.word_wrap = True
    p_l = tf_l.paragraphs[0]
    p_l.text = "03. TF-IDF 가중치의 시대정신"
    set_font(p_l.runs[0], "Gmarket Sans", 20, bold=True, color_rgb=TXT_DARK)
    p_l2 = tf_l.add_paragraph()
    p_l2.text = (
        "\n■ 트렌드 이동: 주식투자, 투자전략, 부자되는법 등 지배\n"
        "■ 독서의 도구화: 책이 단순 교양이 아닌 개인이 경제 불황에 "
        "살아남기 위한 실용적인 생존 학습 채널로 전환됨.\n"
        "■ 제언: 이론 중심 서적을 배제하고 즉시 엑셀 가계부나 "
        "포트폴리오에 대입해 실행할 수 있는 행동 지향 매뉴얼 "
        "기획이 압도적 베스트셀러 진입 확률을 확보함."
    )
    set_font(p_l2.runs[0], "Nanum Gothic", 13, color_rgb=RGBColor(60, 60, 60))
    
    add_card_shape(slide29, Inches(6.8), Inches(1.8), Inches(5.6), Inches(4.8))
    tb_r = slide29.shapes.add_textbox(Inches(7.0), Inches(2.0), Inches(5.2), Inches(4.4))
    tf_r = tb_r.text_frame
    tf_r.word_wrap = True
    p_r = tf_r.paragraphs[0]
    p_r.text = "04. 리뷰 플라이휠 비선형 가속성"
    set_font(p_r.runs[0], "Gmarket Sans", 20, bold=True, color_rgb=TXT_DARK)
    p_r2 = tf_r.add_paragraph()
    p_r2.text = (
        "\n■ 수치 검증: 리뷰 10건이하 지수 3천점 vs 100건돌파시 "
        "지수 3만점대로 수직 급상승 (10배 비선형 급등)\n"
        "■ 추진 과제: 도서정가제 고정 10%로 가격 혜택이 불가하므로, "
        "출간 초기 2 ~ 4주 골든타임 이내에 집중 마케팅 화력을 "
        "동원해 양질 리뷰 100건 임계값을 조기 확보하여 플랫폼 "
        "알고리즘 상위 노출 기회를 독점해야 함."
    )
    set_font(p_r2.runs[0], "Nanum Gothic", 13, color_rgb=RGBColor(60, 60, 60))
    
    note29 = (
        "종합 인사이트 3과 4입니다. 자연어 마이닝이 규명한 시대적 코드는 교양 담론이 아니라 '생존을 위한 도구적 학습'입니다. "
        "행동 지침을 곧바로 실행 가능한 가계부나 연금 시뮬레이터 번들과 매뉴얼 형태로 서적을 설계해야 기획의 당위성을 "
        "확보합니다. 또한, 리뷰가 100건을 넘기는 시점부터 판매량이 10배 급등하는 플라이휠 효과가 실증된바, 출간 "
        "초기 2주 내에 대량의 리뷰를 축적해 알고리즘 우대를 선점하는 타임 마케팅이 베스트셀러 기획의 절대적인 "
        "생존 공식입니다."
    )
    slide29.notes_slide.notes_text_frame.text = note29

    # -------------------------------------------------------------
    # Slide 30. [간지] PART 7. 전략적 제언 및 결론
    # -------------------------------------------------------------
    slide30 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide30, BG_DARK)
    
    add_card_shape(slide30, Inches(1.0), Inches(2.0), Inches(11.33), Inches(3.5), fill_rgb=RGBColor(40, 35, 36), border_rgb=COLOR_GOLD)
    tb = slide30.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(10.33), Inches(2.5))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "PART 7. 전략적 제언 및 성공 프레임워크"
    set_font(p.runs[0], "Gmarket Sans", 36, bold=True, color_rgb=TXT_LIGHT)
    p2 = tf.add_paragraph()
    p2.text = "베스트셀러 흥행을 주도하기 위한 5대 기획 프레임워크 제언"
    set_font(p2.runs[0], "Nanum Gothic", 18, color_rgb=COLOR_GOLD)
    
    note30 = (
        "마지막 일곱 번째 대분류인 전략적 제언 및 결론 파트입니다. 수치로 입증한 4대인사이트와 요인들을 융합하여 "
        "실무 기획자가 신작 출간 즉시 실행에 옮길 수 있는 '5대 흥행 성공 공식 프레임워크'와 리스크 헤징을 위한 "
        "하이브리드 유통 모델 제안, 그리고 최종 요약 보고를 드리며 질의응답을 열도록 하겠습니다. 끝까지 경청해 "
        "주셔서 대단히 감사합니다."
    )
    slide30.notes_slide.notes_text_frame.text = note30

    # -------------------------------------------------------------
    # Slide 31. 베스트셀러 흥행 5대 프레임워크 (인포그래픽 카드)
    # -------------------------------------------------------------
    slide31 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide31, BG_LIGHT)
    add_title(slide31, "7.1 베스트셀러 흥행 기획 5대 실천 프레임워크")
    
    # 5개 가로 카드 배치
    frameworks = [
        ("01. 타겟 불안 저격", "은퇴 통장 설계, 1억 모으기 등 원초적 생존 불안 요소를 책의 핵심 타이틀로 매핑"),
        ("02. 행동 매뉴얼 구성", "추상 이론을 배제하고 즉시 행동할 수 있는 도구, 템플릿, 실전 워크북 번들링"),
        ("03. 리뷰 플라이휠 조직", "출간 2 ~ 4주 이내에 서평단 100명을 가동해 초기 리뷰 100건 임계점 돌파 가속"),
        ("04. 상위 출판사 파트너", "유통망 장악 및 초기 마케팅 자금력이 확보된 10대 브랜드와의 기획 계약"),
        ("05. 구독 플랫폼 선검증", "크레마클럽 등 전자책 채널에 반응 지표를 선조사한 뒤 종이책 소장용으로 2차 출판")
    ]
    for i, (fw_t, fw_d) in enumerate(frameworks):
        y = Inches(1.6 + (i * 1.05))
        add_card_shape(slide31, Inches(0.8), y, Inches(11.73), Inches(0.95))
        
        tb = slide31.shapes.add_textbox(Inches(1.0), y + Inches(0.1), Inches(11.33), Inches(0.75))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"★ {fw_t}  |  {fw_d}"
        set_font(p.runs[0], "Nanum Gothic", 14, bold=True, color_rgb=TXT_DARK)
        
    note31 = (
        "5대 흥행 실천 프레임워크 상세 지침입니다. 첫째, 도서의 핵심 테마는 타겟층의 금융 및 부동산, 미래 노후에 "
        "대한 원초적 생존 불안을 저격해야 마찰을 뚫습니다. 둘째, 바로 입력하고 검증할 수 있는 재테크 가이드 템플릿을 "
        "포함하십시오. 셋째, 출간 초기 골든타임 내에 100건의 평가 글을 조직하여 플랫폼 알고리즘 추천을 획득합니다. "
        "넷째, 초기 푸시와 배포력이 보증된 대형 출판사와 협업하며, 다섯째, 전자책 구독 플랫폼의 지표 추이를 보고 "
        "반응이 좋은 경우에만 종이책을 발행해 리스크를 최소화하는 의사결정 프레임워크를 수립해야 흥행 승률을 올립니다."
    )
    slide31.notes_slide.notes_text_frame.text = note31

    # -------------------------------------------------------------
    # Slide 32. 하이브리드 리스크 헤징 모델 (비교 그리드)
    # -------------------------------------------------------------
    slide32 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide32, BG_LIGHT)
    add_title(slide32, "7.2 리스크 헤징을 위한 하이브리드 출판 모델")
    
    # 2열 카드 배치
    add_card_shape(slide32, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8))
    tb_l = slide32.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(5.2), Inches(4.4))
    tf_l = tb_l.text_frame
    tf_l.word_wrap = True
    p_l = tf_l.paragraphs[0]
    p_l.text = "1단계: 디지털 구독 선검증"
    set_font(p_l.runs[0], "Gmarket Sans", 20, bold=True, color_rgb=TXT_DARK)
    p_l2 = tf_l.add_paragraph()
    p_l2.text = (
        "\n■ 채널: 크레마클럽 등 전자책 정기 구독 채널 선출간\n"
        "■ 검증 지표: 독서 시작율, 평균 완독률, 평점 추이 분석\n"
        "■ 장점: 인쇄 및 유통 물류 초기 비용을 90% 이상 절감하면서 "
        "실제 시장 반응과 독자 리뷰 텍스트를 가장 빠르게 확보."
    )
    set_font(p_l2.runs[0], "Nanum Gothic", 14, color_rgb=RGBColor(60, 60, 60))
    
    add_card_shape(slide32, Inches(6.8), Inches(1.8), Inches(5.6), Inches(4.8))
    tb_r = slide32.shapes.add_textbox(Inches(7.0), Inches(2.0), Inches(5.2), Inches(4.4))
    tf_r = tb_r.text_frame
    tf_r.word_wrap = True
    p_r = tf_r.paragraphs[0]
    p_r.text = "2단계: 프리미엄 소장 단행본화"
    set_font(p_r.runs[0], "Gmarket Sans", 20, bold=True, color_rgb=TXT_DARK)
    p_r2 = tf_r.add_paragraph()
    p_r2.text = (
        "\n■ 채널: 오프라인 대형 서점 및 예스24 메인 지대 노출\n"
        "■ 고도화 작업: 선검증 리뷰 데이터 내 불만 보완 및 증면\n"
        "■ 장점: 디지털 독자 평을 마케팅 서평단으로 즉시 치환하고 "
        "흥행 승률이 90% 이상 담보된 상태에서 종이책 물류 투하."
    )
    set_font(p_r2.runs[0], "Nanum Gothic", 14, color_rgb=RGBColor(60, 60, 60))
    
    note32 = (
        "출판 리스크를 제어하기 위한 하이브리드 의사결정 모델입니다. 최근 경제경영 도서의 트렌드는 종이책을 무작정 "
        "찍어 매대에 올리는 모험을 지양합니다. 1단계로 크레마클럽 등 디지털 플랫폼에 epub 형태로 선출시하여 조회수와 "
        "완독률을 확인합니다. 데이터를 통해 상업적 파급력이 입증되면 독자들의 피드백을 책의 개정판에 적용해 완성도를 "
        "더욱 올린 후, 프리미엄 종이책 소장본으로 오프라인 서점에 2차 배포함으로써 재고 리스크를 혁신적으로 줄입니다."
    )
    slide32.notes_slide.notes_text_frame.text = note32

    # -------------------------------------------------------------
    # Slide 33. 종합 결론 및 향후 전망
    # -------------------------------------------------------------
    slide33 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide33, BG_LIGHT)
    add_title(slide33, "7.3 종합 결론 및 데이터 기반 기획 전망")
    
    # 2열 비교 카드 배치
    add_card_shape(slide33, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8))
    tb_l = slide33.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(5.2), Inches(4.4))
    tf_l = tb_l.text_frame
    tf_l.word_wrap = True
    p_l = tf_l.paragraphs[0]
    p_l.text = "데이터 기반 의사결정 체계화"
    set_font(p_l.runs[0], "Gmarket Sans", 20, bold=True, color_rgb=TXT_DARK)
    p_l2 = tf_l.add_paragraph()
    p_l2.text = (
        "\n■ 기획 체계화: 평점 쏠림 필터링을 거친 객관성 확보\n"
        "■ 리뷰건수 가속화: 비선형 시너지 구간의 정량적 제어\n"
        "■ 트렌드 반영: 형태소 분석을 통한 실질적 시대 불안 저격\n\n"
        "전통적인 경험과 추정에 기초하던 도서 기획을 정량 통계에 근거한 의사결정 모델로 표준화해야 합니다."
    )
    set_font(p_l2.runs[0], "Nanum Gothic", 14, color_rgb=RGBColor(60, 60, 60))
    
    add_card_shape(slide33, Inches(6.8), Inches(1.8), Inches(5.6), Inches(4.8))
    tb_r = slide33.shapes.add_textbox(Inches(7.0), Inches(2.0), Inches(5.2), Inches(4.4))
    tf_r = tb_r.text_frame
    tf_r.word_wrap = True
    p_r = tf_r.paragraphs[0]
    p_r.text = "향후 독서 소비 시장 전망"
    set_font(p_r.runs[0], "Gmarket Sans", 20, bold=True, color_rgb=TXT_DARK)
    p_r2 = tf_r.add_paragraph()
    p_r2.text = (
        "\n■ 지식 도구화 가속: 재테크 및 생존 테크닉 지속 주도\n"
        "■ 소장 가치 극대화: 단순 대여는 구독으로, 소장은 고급화로\n"
        "■ 플랫폼 종속 심화: 추천 알고리즘 획득이 마케팅의 성패 결정\n\n"
        "앞으로의 독자들은 단순 지식 획득이 아닌 직접 가치 연마를 원하며, 이에 적응하는 콘텐츠만이 베스트셀러를 승계합니다."
    )
    set_font(p_r2.runs[0], "Nanum Gothic", 14, color_rgb=RGBColor(60, 60, 60))
    
    note33 = (
        "오늘 발표의 종합 결론이자 미래 도서 기획의 나침반 전망입니다. 본 분석은 감각의 영역에만 머물던 출판 행위를 "
        "정량 데이터를 통해 설명 가능한 비즈니스로 고도화했습니다. 앞으로 독서 소비는 더욱 '생존의 지식 도구화'로 "
        "가속될 것이며 소장을 원하게 만드는 고급 양장본과 플랫폼 정기 대여 시장으로 양극화될 것입니다. 이 시대의 흐름을 "
        "포착하여 데이터 기반 기획 프레임워크를 적용하는 기업만이 베스트셀러 왕좌를 오랫동안 점유할 것입니다."
    )
    slide33.notes_slide.notes_text_frame.text = note33

    # -------------------------------------------------------------
    # Slide 34. Q&A 및 감사 인사 (마무리 Slide)
    # -------------------------------------------------------------
    slide34 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide34, BG_DARK)
    
    add_card_shape(slide34, Inches(0.8), Inches(1.8), Inches(11.73), Inches(3.8), fill_rgb=RGBColor(45, 40, 41), border_rgb=COLOR_GOLD)
    
    title_box = slide34.shapes.add_textbox(Inches(1.2), Inches(2.2), Inches(11.0), Inches(2.0))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "경청해 주셔서 감사합니다."
    set_font(p.runs[0], "Gmarket Sans", 42, bold=True, color_rgb=TXT_LIGHT)
    
    p2 = tf.add_paragraph()
    p2.text = "\n질의응답(Q&A) 시간이오니 편하게 질문해 주시기 바랍니다."
    set_font(p2.runs[0], "Nanum Gothic", 18, color_rgb=COLOR_GOLD)
    
    p3 = tf.add_paragraph()
    p3.text = "\n* 상세 보고서 본문: eda_report.md  |  * 분석용 데이터 대시보드: bestsellers_dashboard.xlsx"
    set_font(p3.runs[0], "Nanum Gothic", 13, color_rgb=TXT_LIGHT)
    
    note34 = (
        "이상으로 준비한 예스24 베스트셀러 고도화 EDA 종합 보고 발표를 모두 마치도록 하겠습니다. 데이터 통계에 대한 "
        "정밀한 추가 수치와 수식 관계는 제공해 드린 상세 마크다운 파일 [eda_report.md] 및 동적 엑셀 분석 도구인 "
        "[bestsellers_dashboard.xlsx] 파일을 통해 한글 주석과 함께 즉시 대입해 검증하실 수 있습니다. 혹시 오늘 다룬 "
        "일변량 가격/평점 분포의 편향성이나 리뷰 건수 임계점 상관성, 혹은 5대 성공 프레임워크 실천 방안에 관하여 "
        "세부적인 질의나 궁금하신 사항이 있으시다면 말씀해 주십시오. 상세히 설명해 올리겠습니다. 감사합니다."
    )
    slide34.notes_slide.notes_text_frame.text = note34

    # 3. PPTX 파일 저장
    prs.save(output_path)
    print(f"\n최종 성공: PPTX 프레젠테이션 빌드 완료! 파일 경로: {output_path}")

if __name__ == "__main__":
    build_presentation()
