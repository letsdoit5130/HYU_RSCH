"""
수출 EDA 대시보드 및 마진/포트폴리오 시장개척 전략 PowerPoint (.pptx) v2 생성 스크립트.
"""

import os
import sys
import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

sys.stdout.reconfigure(encoding='utf-8')

PPTX_OUTPUT_PATH = 'BIZ-laver/reports/HaeYu_Laver_Export_Market_Strategy_Presentation_v2.pptx'
IMAGE_DIR = 'BIZ-laver/images'

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

COLOR_BG_DARK = RGBColor(15, 23, 42)
COLOR_NAVY = RGBColor(26, 82, 118)
COLOR_ACCENT_BLUE = RGBColor(52, 152, 219)
COLOR_TEXT_MAIN = RGBColor(30, 41, 59)
COLOR_CARD_BG = RGBColor(248, 250, 252)

blank_layout = prs.slide_layouts[6]

def add_slide_header(slide, title_text, category_text="해유 김 수출 심화 EDA & 1인 무역회사 시장개척"):
    header_box = slide.shapes.add_shape(pptx.enum.shapes.MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.1))
    header_box.fill.solid()
    header_box.fill.fore_color.rgb = COLOR_NAVY
    header_box.line.color.rgb = COLOR_NAVY
    
    tf = header_box.text_frame
    tf.margin_left = Inches(0.5)
    tf.margin_top = Inches(0.15)
    
    p0 = tf.paragraphs[0]
    p0.text = category_text.upper()
    p0.font.size = Pt(10)
    p0.font.bold = True
    p0.font.color.rgb = COLOR_ACCENT_BLUE
    p0.font.name = '맑은 고딕'
    
    p1 = tf.add_paragraph()
    p1.text = title_text
    p1.font.size = Pt(20)
    p1.font.bold = True
    p1.font.color.rgb = RGBColor(255, 255, 255)
    p1.font.name = '맑은 고딕'

# Slide 1
slide1 = prs.slides.add_slide(blank_layout)
bg1 = slide1.shapes.add_shape(pptx.enum.shapes.MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
bg1.fill.solid()
bg1.fill.fore_color.rgb = COLOR_BG_DARK
bg1.line.color.rgb = COLOR_BG_DARK

tf1 = slide1.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.333), Inches(3.5)).text_frame
p = tf1.paragraphs[0]
p.text = "해유 김 수출 EDA & 심화 시장개척 전략 (v2)"
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)
p.font.name = '맑은 고딕'

p_sub = tf1.add_paragraph()
p_sub.text = "1인 무역회사 관점의 4분면 시장 포트폴리오, 단가 마진 구간 및 잠재 파트너 개척안"
p_sub.font.size = Pt(18)
p_sub.font.color.rgb = COLOR_ACCENT_BLUE
p_sub.font.name = '맑은 고딕'
p_sub.space_before = Pt(12)

# Slide 2
slide2 = prs.slides.add_slide(blank_layout)
add_slide_header(slide2, "글로벌 40개국 4분면 포트폴리오 매트릭스 (성장률 x 단가 마진)")

img14 = os.path.join(IMAGE_DIR, '14_advanced_market_portfolio_matrix.png')
if os.path.exists(img14):
    slide2.shapes.add_picture(img14, Inches(0.6), Inches(1.4), Inches(7.2), Inches(5.4))

card_box = slide2.shapes.add_shape(pptx.enum.shapes.MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.0), Inches(1.4), Inches(4.7), Inches(5.4))
card_box.fill.solid()
card_box.fill.fore_color.rgb = COLOR_CARD_BG
card_box.line.color.rgb = COLOR_NAVY

tf_card = card_box.text_frame
tf_card.margin_left = Inches(0.3)
tf_card.margin_top = Inches(0.3)
tf_card.word_wrap = True

p = tf_card.paragraphs[0]
p.text = "★ Star Market (1순위 개척)"
p.font.size = Pt(15)
p.font.bold = True
p.font.color.rgb = RGBColor(46, 134, 193)

bullet_points = [
    "UAE: 5년 성장 +187.5%, 단가 $25.06/kg (할랄 프리미엄 스낵)",
    "폴란드: 5년 성장 +121.2%, 단가 $24.01/kg (동유럽 리테일 허브)",
    "콜롬비아: 5년 성장 +156.7%, 단가 $27.47/kg (남미 고마진)",
    "튀르키예: 5년 성장 +254.5%, 단가 $28.01/kg (초고마진 틈새)",
    "",
    "🚀 1인 무역상 액션 플랜:",
    "• 조미김 완제품(HS 200899) 스낵형 제품 100% 집중",
    "• $20~$30/kg 이상 소비자가격 형성을 위한 고마진 포지셔닝",
    "• LCL 소량 해상 운송 & 현지 크로스보더 E-Commerce 직입점"
]

for bp in bullet_points:
    p_b = tf_card.add_paragraph()
    p_b.text = bp
    p_b.font.size = Pt(11)
    p_b.font.color.rgb = COLOR_TEXT_MAIN
    p_b.font.name = '맑은 고딕'

# Slide 3
slide3 = prs.slides.add_slide(blank_layout)
add_slide_header(slide3, "수출 단가 5대 마진 구간 및 국가별 단가 변동성 (CV%) 리스크")

img13 = os.path.join(IMAGE_DIR, '13_advanced_price_bracket_distribution.png')
img15 = os.path.join(IMAGE_DIR, '15_advanced_seasoned_vs_raw_monthly_trend.png')

if os.path.exists(img13):
    slide3.shapes.add_picture(img13, Inches(0.6), Inches(1.4), Inches(5.8), Inches(5.4))

if os.path.exists(img15):
    slide3.shapes.add_picture(img15, Inches(6.7), Inches(1.4), Inches(6.0), Inches(5.4))

prs.save(PPTX_OUTPUT_PATH)
print(f"PowerPoint 발표 덱이 {PPTX_OUTPUT_PATH}에 성공적으로 재생성 저장되었습니다.")
