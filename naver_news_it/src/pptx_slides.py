"""
naver_news_it PPTX 슬라이드 발표자료 생성 스크립트

이 모듈은 python-pptx 라이브러리를 사용하여 수집된 데이터와 
시각화 분석 결과를 슬라이드 장표로 구현한 발표자료(.pptx)를 자동 작성합니다.

작성일: 2026-07-19
"""

import os
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt

def build_pptx_slides():
    print("[PPTX Slides] 발표 슬라이드 생성을 시작합니다...")
    
    current_dir = os.path.dirname(os.path.abspath(__file__))
    data_path = os.path.abspath(os.path.join(current_dir, "..", "data", "raw_data.csv"))
    output_path = os.path.abspath(os.path.join(current_dir, "..", "docs", "presentation.pptx"))
    
    if not os.path.exists(data_path):
        print(f"[PPTX Slides] [ERROR] 원시 데이터가 존재하지 않습니다: {data_path}")
        return
        
    df = pd.read_csv(data_path)
    
    # PPTX Presentation 생성
    prs = Presentation()
    
    # 슬라이드 마스터 레이아웃: 0번은 표지, 1번은 본문 레이아웃
    slide_layout_title = prs.slide_layouts[0]
    slide_layout_content = prs.slide_layouts[1]
    
    # 1. 표지 슬라이드 추가
    slide = prs.slides.add_slide(slide_layout_title)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "네이버 IT 뉴스 데이터 분석"
    subtitle.text = f"수집 데이터 기반 언론사 분포 및 리포팅\n분석일자: {df['date'].iloc[0] if len(df) > 0 else 'N/A'} (총 {len(df)}건 분석)"
    
    # 2. 본문 슬라이드 (데이터 개요)
    slide2 = prs.slides.add_slide(slide_layout_content)
    shapes = slide2.shapes
    title_shape = shapes.title
    title_shape.text = "1. 뉴스 데이터 수집 및 구성 개요"
    
    body_shape = shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.text = f"네이버 뉴스 IT/과학 섹션에서 총 {len(df)}건의 기사 정보를 수집하였습니다."
    
    p = tf.add_paragraph()
    p.text = f"- 수집 필드: 기사 제목(title), 요약(summary), 언론사(publisher), 기사 링크(link)"
    p.level = 1
    
    p2 = tf.add_paragraph()
    p2.text = f"- 데이터 보존: 중복 제거(제목 기준)를 완료하여 원천 데이터 정합성 확보"
    p2.level = 1
    
    # 3. 본문 슬라이드 (시각화 차트)
    slide3 = prs.slides.add_slide(slide_layout_content)
    title_shape3 = slide3.shapes.title
    title_shape3.text = "2. 언론사별 기사 분포 시각화"
    
    chart_path = os.path.abspath(os.path.join(current_dir, "..", "images", "summary_chart.png"))
    if os.path.exists(chart_path):
        left = Inches(0.8)
        top = Inches(2.0)
        width = Inches(5.8)
        height = Inches(4.2)
        slide3.shapes.add_picture(chart_path, left, top, width, height)
        
        # 텍스트 상자 추가
        txBox = slide3.shapes.add_textbox(Inches(6.8), Inches(2.0), Inches(3.0), Inches(4.2))
        tf_box = txBox.text_frame
        tf_box.word_wrap = True
        tf_box.text = "시사점 및 분석 해석:"
        
        p = tf_box.add_paragraph()
        p.text = "- 특정 미디어(종합지/IT 전문지)가 뉴스 생산의 높은 점유율을 차지하고 있습니다."
        p.level = 0
        
        p_sub = tf_box.add_paragraph()
        p_sub.text = "- 이는 IT 분야 기사의 생산 집중도가 특정 언론사에 쏠려 있음을 시사합니다."
        p_sub.level = 0
    else:
        # 차트 부재 시 텍스트 프레임으로 대체
        body_shape3 = slide3.shapes.placeholders[1]
        body_shape3.text = "[요약 차트 이미지 생략됨]"
        
    prs.save(output_path)
    print(f"[PPTX Slides] PPTX 슬라이드 저장 완료: {output_path}")
    return output_path

if __name__ == "__main__":
    build_pptx_slides()
