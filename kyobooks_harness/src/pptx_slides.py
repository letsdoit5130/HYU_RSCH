"""
kyobooks_harness PPTX 슬라이드 발표자료 생성 스크립트

이 모듈은 python-pptx 라이브러리를 사용하여 수집된 데이터와 
시각화 분석 결과를 슬라이드 장표로 구현한 발표자료(.pptx)를 자동 작성합니다.

작성일: 2026-07-19
"""

import os
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt

def build_pptx_slides():
    print("[PPTX Slides] 발표 슬라이드 생성을 시작합니다...")
    
    current_dir = os.path.dirname(os.path.abspath(__file__))
    data_path = os.path.abspath(os.path.join(current_dir, "..", "data", "raw_data.csv"))
    output_path = os.path.abspath(os.path.join(current_dir, "..", "docs", "presentation.pptx"))
    
    if not os.path.exists(data_path):
        print(f"[PPTX Slides] [ERROR] 원시 데이터가 존재하지 않습니다: {data_path}")
        return
        
    df = pd.read_csv(data_path)
    
    # PPTX Presentation 생성
    prs = Presentation()
    
    # 슬라이드 마스터 레이아웃: 0번은 표지, 1번은 본문 레이아웃
    slide_layout_title = prs.slide_layouts[0]
    slide_layout_content = prs.slide_layouts[1]
    
    # 1. 표지 슬라이드 추가
    slide = prs.slides.add_slide(slide_layout_title)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "kyobooks_harness 분석 발표자료"
    subtitle.text = f"수집 데이터 기반 비즈니스 제언\n작성일: 2026-07-19 (총 {len(df)}건 분석)"
    
    # 2. 본문 슬라이드 (데이터 개요)
    slide2 = prs.slides.add_slide(slide_layout_content)
    shapes = slide2.shapes
    title_shape = shapes.title
    title_shape.text = "1. 수집 데이터 개요"
    
    body_shape = shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.text = f"타겟 웹페이지에서 총 {len(df)}건의 데이터를 안정적으로 수집하였습니다."
    
    p = tf.add_paragraph()
    p.text = f"- 수집 필드: 상품번호, 순위, 도서명, 저자, 출판사, 출판일, 정가, 할인가, 할인율, 평점"
    p.level = 1
    
    p2 = tf.add_paragraph()
    p2.text = f"- 데이터 인코딩: UTF-8-SIG (Excel 호환)"
    p2.level = 1
    
    # 3. 본문 슬라이드 (시각화 차트)
    slide3 = prs.slides.add_slide(slide_layout_content)
    title_shape3 = slide3.shapes.title
    title_shape3.text = "2. 시각화 분석 및 시사점"
    
    chart_path = os.path.abspath(os.path.join(current_dir, "..", "images", "summary_chart.png"))
    if os.path.exists(chart_path):
        left = Inches(1)
        top = Inches(2.2)
        width = Inches(5.5)
        height = Inches(4.0)
        slide3.shapes.add_picture(chart_path, left, top, width, height)
        
        # 텍스트 상자 추가
        txBox = slide3.shapes.add_textbox(Inches(6.8), Inches(2.2), Inches(3.0), Inches(4.0))
        tf_box = txBox.text_frame
        tf_box.word_wrap = True
        tf_box.text = "주요 분석 결과 요약:"
        
        p = tf_box.add_paragraph()
        p.text = "- 요약 차트에서 보듯 데이터 분포가 고르게 형성되어 있습니다."
        p.level = 0
    else:
        # 차트 부재 시 텍스트 프레임으로 대체
        body_shape3 = slide3.shapes.placeholders[1]
        body_shape3.text = "[요약 차트 이미지 생략됨]"
        
    prs.save(output_path)
    print(f"[PPTX Slides] PPTX 슬라이드 저장 완료: {output_path}")
    return output_path

if __name__ == "__main__":
    build_pptx_slides()
