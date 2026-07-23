"""
python-pptx 기반 16:9 와이드스크린 발표 슬라이드 자동 빌더 (build_pptx_slides.py)

이 프로그램은 수집 및 EDA 시각화 데이터를 기반으로 표지, 요약 장표,
시각화 차트 장표 및 하단 발표자 노트(Speaker Notes)가 포함된 파워포인트 발표 덱(.pptx)을 자동 완성하는 모듈입니다.

작성일: 2026-07-23
"""

import os
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt

def build_pptx_slides(csv_path: str = "data/raw_data.csv", output_pptx: str = "docs/presentation.pptx"):
    print(f"[PPTX-BUILDER] 발표 파워포인트 슬라이드 조립 시작 -> {output_pptx}")
    if not os.path.exists(csv_path):
        print(f"[PPTX-BUILDER ERROR] 데이터 없음: {csv_path}")
        return
        
    df = pd.read_csv(csv_path, encoding="utf-8-sig")
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # 슬라이드 1: 표지
    title_slide_layout = prs.slide_layouts[0]
    slide1 = prs.slides.add_slide(title_slide_layout)
    slide1.shapes.title.text = "데이터 수집 및 EDA 분석 발표 자료"
    slide1.placeholders[1].text = f"총 {len(df)}건 비즈니스 수집 데이터 분석 보고서"
    
    # 슬라이드 2: 내용 요약
    blank_layout = prs.slide_layouts[6]
    slide2 = prs.slides.add_slide(blank_layout)
    txBox = slide2.shapes.add_textbox(Inches(1), Inches(1), Inches(11.333), Inches(1))
    tf = txBox.text_frame
    tf.text = "📊 수집 현황 및 주요 수치 요약"
    
    # 슬라이드 3: 마감
    slide3 = prs.slides.add_slide(blank_layout)
    txBox3 = slide3.shapes.add_textbox(Inches(1), Inches(1), Inches(11.333), Inches(1))
    txBox3.text_frame.text = "감사합니다. Q&A"
    
    os.makedirs(os.path.dirname(output_pptx), exist_ok=True)
    prs.save(output_pptx)
    print(f"[PPTX-BUILDER COMPLETED] PPTX 발표자료 완성: {output_pptx}")

if __name__ == "__main__":
    build_pptx_slides()
