"""
통합 마스터 시장개척 발표 슬라이드 PowerPoint (.pptx) 파일 생성 스크립트.

기능:
1. 16:9 와이드스크린 레이아웃 및 Midnight Executive 디자인 시스템 적용
2. Executive Summary, 4분면 시장 포트폴리오, 조미김 Top 10 바이어 파트너십, 실전 4단계 개척 프로토콜 슬라이드 구성
3. BIZ-laver/reports/HaeYu_Laver_Export_Master_Market_Expansion_Deck.pptx 생성
"""

import os
import sys
import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

sys.stdout.reconfigure(encoding='utf-8')

PPTX_OUTPUT_PATH = 'BIZ-laver/reports/HaeYu_Laver_Export_Master_Market_Expansion_Deck.pptx'
IMAGE_DIR = 'BIZ-laver/images'

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

COLOR_BG_DARK = RGBColor(15, 23, 42)
COLOR_NAVY = RGBColor(26, 82, 118)
COLOR_ACCENT_BLUE = RGBColor(52, 152, 219)
COLOR_TEXT_MAIN = RGBColor(30, 41, 59)
COLOR_CARD_BG = RGBColor(248, 250, 252)

blank_layout = prs.slide_layouts[6]

def add_slide_header(slide, title_text, category_text="해유 김 수출 종합 EDA & 1인 무역회사 시장개척"):
    header_box = slide.shapes.add_shape(pptx.enum.shapes.MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.1))
    header_box.fill.solid()
    header_box.fill.fore_color.rgb = COLOR_NAVY
    header_box.line.color.rgb = COLOR_NAVY
    
    tf = header_box.text_frame
    tf.margin_left = Inches(0.5)
    tf.margin_top = Inches(0.15)
    
    p0 = tf.paragraphs[0]
    p0.text = category_text.upper()
    p0.font.size = Pt(10)
    p0.font.bold = True
    p0.font.color.rgb = COLOR_ACCENT_BLUE
    p0.font.name = '맑은 고딕'
    
    p1 = tf.add_paragraph()
    p1.text = title_text
    p1.font.size = Pt(20)
    p1.font.bold = True
    p1.font.color.rgb = RGBColor(255, 255, 255)
    p1.font.name = '맑은 고딕'

# Slide 1: Cover (Dark Theme)
slide1 = prs.slides.add_slide(blank_layout)
bg1 = slide1.shapes.add_shape(pptx.enum.shapes.MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
bg1.fill.solid()
bg1.fill.fore_color.rgb = COLOR_BG_DARK
bg1.line.color.rgb = COLOR_BG_DARK

tf1 = slide1.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.333), Inches(3.5)).text_frame
p = tf1.paragraphs[0]
p.text = "[마스터 전략] 글로벌 김 수출 EDA & 1인 무역회사 시장개척"
p.font.size = Pt(34)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)
p.font.name = '맑은 고딕'

p_sub = tf1.add_paragraph()
p_sub.text = "Star Market 4개국, 조미김 고마진 포지셔닝 및 4단계 실전 개척 프로토콜"
p_sub.font.size = Pt(18)
p_sub.font.color.rgb = COLOR_ACCENT_BLUE
p_sub.font.name = '맑은 고딕'
p_sub.space_before = Pt(12)

# Slide 2: 4분면 시장 포트폴리오
slide2 = prs.slides.add_slide(blank_layout)
add_slide_header(slide2, "글로벌 40개국 4분면 포트폴리오 매트릭스 (Star Market 4개국)")

img14 = os.path.join(IMAGE_DIR, '14_advanced_market_portfolio_matrix.png')
if os.path.exists(img14):
    slide2.shapes.add_picture(img14, Inches(0.6), Inches(1.4), Inches(7.2), Inches(5.4))

card_box = slide2.shapes.add_shape(pptx.enum.shapes.MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.0), Inches(1.4), Inches(4.7), Inches(5.4))
card_box.fill.solid()
card_box.fill.fore_color.rgb = COLOR_CARD_BG
card_box.line.color.rgb = COLOR_NAVY

tf_card = card_box.text_frame
tf_card.margin_left = Inches(0.3)
tf_card.margin_top = Inches(0.3)
tf_card.word_wrap = True

p = tf_card.paragraphs[0]
p.text = "★ Star Market 4개국 (최우선 개척)"
p.font.size = Pt(15)
p.font.bold = True
p.font.color.rgb = RGBColor(46, 134, 193)

bullet_points = [
    "1. UAE: 5년 성장 +187.5%, 단가 $25.06/kg (할랄 K-스낵)",
    "2. 폴란드: 5년 성장 +121.2%, 단가 $24.01/kg (동유럽 Hub)",
    "3. 콜롬비아: 5년 성장 +156.7%, 단가 $27.47/kg (남미 고마진)",
    "4. 튀르키예: 5년 성장 +254.5%, 단가 $28.01/kg (초고마진 틈새)",
    "",
    "💡 1인 무역상 전략 핵심:",
    "• 조미김 완제품(HS 200899) 스낵에 100% 사격 집중",
    "• $20~$30/kg 이상 소비자가격을 형성하는 프리미엄 포지셔닝",
    "• 초기 LCL 소량 수송 및 현지 E-Commerce 직입점"
]

for bp in bullet_points:
    p_b = tf_card.add_paragraph()
    p_b.text = bp
    p_b.font.size = Pt(11)
    p_b.font.color.rgb = COLOR_TEXT_MAIN
    p_b.font.name = '맑은 고딕'

# Slide 3: 조미김 Top 10 바이어 파트너십
slide3 = prs.slides.add_slide(blank_layout)
add_slide_header(slide3, "조미김 Top 10 유망 타깃 시장 & 현지 바이어 파트너십")

card_box3 = slide3.shapes.add_shape(pptx.enum.shapes.MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.4), Inches(12.133), Inches(5.4))
card_box3.fill.solid()
card_box3.fill.fore_color.rgb = COLOR_CARD_BG
card_box3.line.color.rgb = COLOR_NAVY

tf3 = card_box3.text_frame
tf3.margin_left = Inches(0.4)
tf3.margin_top = Inches(0.3)

p3 = tf3.paragraphs[0]
p3.text = "🎯 주요 국가별 컨택 가능 잠재 바이어 / 수입상 리스트"
p3.font.size = Pt(16)
p3.font.bold = True
p3.font.color.rgb = COLOR_NAVY

top10_list = [
    "• USA: Assi Rhee Bros, H-Mart, Harvesko, Weee!",
    "• Japan: E-Mart Japan, Gyomu Super, CJ CheilJedang Japan",
    "• Russian Fed.: X5 Retail Group, Magnit, Koros Co.",
    "• Canada: T&T Supermarket, PAT Mart, Galleria Supermarket",
    "• Poland: Kuchnie Świata, Asian House Poland, Allegro Sellers",
    "• UAE: Choithrams, Lulu Hypermarket, Kibsons, Al Maya Group",
    "• Kazakhstan: Magnum Cash & Carry, Shin-Line, Small Supermarket",
    "• Türkiye: Macrocenter, Gurme Park, Happy Center"
]

for item in top10_list:
    p_item = tf3.add_paragraph()
    p_item.text = item
    p_item.font.size = Pt(13)
    p_item.font.color.rgb = COLOR_TEXT_MAIN
    p_item.font.name = '맑은 고딕'
    p_item.space_before = Pt(6)

# Slide 4: 1인 무역회사 실전 4단계 프로토콜
slide4 = prs.slides.add_slide(blank_layout)
add_slide_header(slide4, "1인 무역회사 실전 4단계 개척 프로토콜")

steps = [
    ("Step 1", "고마진 포지셔닝", "조미김 완제품 스낵\n지퍼백 파우치 소포장\n비건 & 할랄 인증"),
    ("Step 2", "KOTRA 바이어 발굴", "무역관 지사화 사업\nEC21, Kompass 활용\n샘플 무료 항공 배송"),
    ("Step 3", "LCL 물류 & 라벨링", "초기 1~2 Pallet LCL\n해상 운송 재고 최소화\n현지어 영양성분 라벨"),
    ("Step 4", "E-Commerce & 숏폼", "Amazon, Noon, Allegro\n현지 크로스보더 입점\n틱톡 K-Food 숏폼 챌린지")
]

for idx, (s_num, s_title, s_desc) in enumerate(steps):
    left_pos = Inches(0.6 + idx * 3.05)
    s_box = slide4.shapes.add_shape(pptx.enum.shapes.MSO_SHAPE.ROUNDED_RECTANGLE, left_pos, Inches(1.8), Inches(2.8), Inches(4.8))
    s_box.fill.solid()
    s_box.fill.fore_color.rgb = COLOR_CARD_BG
    s_box.line.color.rgb = COLOR_ACCENT_BLUE
    
    tf_s = s_box.text_frame
    tf_s.margin_left = Inches(0.2)
    tf_s.margin_top = Inches(0.3)
    
    p0 = tf_s.paragraphs[0]
    p0.text = s_num
    p0.font.size = Pt(14)
    p0.font.bold = True
    p0.font.color.rgb = COLOR_ACCENT_BLUE
    
    p1 = tf_s.add_paragraph()
    p1.text = s_title
    p1.font.size = Pt(16)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_NAVY
    p1.space_before = Pt(8)
    
    p2 = tf_s.add_paragraph()
    p2.text = s_desc
    p2.font.size = Pt(11)
    p2.font.color.rgb = COLOR_TEXT_MAIN
    p2.space_before = Pt(12)

prs.save(PPTX_OUTPUT_PATH)
print(f"통합 마스터 PowerPoint 발표 덱이 {PPTX_OUTPUT_PATH}에 성공적으로 저장되었습니다.")
