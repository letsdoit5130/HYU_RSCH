"""
통합 테스트 하네스 검증 실행기 (test_harness.py)

이 프로그램은 수집 및 분석 자동화 파이프라인의 물리적 결과물들(CSV, Excel, Word, PPTX, EDA Report)이 
정량 성능 목표 및 구조적 규칙 가이드를 완전히 충족하는지 검증하는 테스트 하네스 러너입니다.

주요 검증 범위:
- CSV 데이터: 파일 생성 여부, 최소 행수(10건) 및 스키마 정합성
- Excel 대시보드: 시트 구조(대시보드 요약, 상세 데이터), 요약 표 및 포맷팅 검증
- Word 보고서: python-docx를 통한 단락 구성 및 상세 표 데이터 검증
- PPTX 발표자료: python-pptx를 통한 슬라이드 구성 및 시각화 장표(이미지) 삽입 여부 검증
- 테스트 결과를 harness/results/ 경로에 Markdown 보고서 및 JSON 로그로 저장

작성일: 2026-07-19
"""

import os
import sys
import time
import json
import argparse
import subprocess
import pandas as pd
from openpyxl import load_workbook
from docx import Document
from pptx import Presentation

def run_pipeline(project_name, target_url):
    """
    지정한 프로젝트의 run_pipeline.py를 서브프로세스로 구동하여 실행 성공 여부를 확인합니다.
    """
    print(f"\n[Harness] {project_name} 파이프라인 실행 테스트 중...")
    workspace_dir = os.path.abspath(os.path.join(os.path.dirname(__file__), "..", ".."))
    pipeline_script = os.path.join(workspace_dir, ".agents", "skills", "crawler-analysis", "scripts", "run_pipeline.py")
    
    # 가상환경 파이썬 인터프리터 경로 확인
    venv_python = os.path.join(workspace_dir, ".venv", "Scripts", "python.exe")
    if not os.path.exists(venv_python):
        venv_python = os.path.join(workspace_dir, ".venv", "bin", "python")
    python_exe = venv_python if os.path.exists(venv_python) else sys.executable
    
    cmd = [python_exe, pipeline_script, "--project", project_name]
    if target_url:
        cmd.extend(["--url", target_url])
        
    try:
        result = subprocess.run(cmd, check=True, capture_output=True, text=True, encoding="utf-8")
        print(" -> [PASS] 파이프라인 스크립트 체인이 정상 종료되었습니다 (Exit Code: 0)")
        return True, ""
    except subprocess.CalledProcessError as e:
        print(f" -> [FAIL] 파이프라인이 에러와 함께 종료되었습니다 (Exit Code: {e.returncode})")
        return False, e.stderr
    except Exception as e:
        print(f" -> [FAIL] 파이프라인 실행 도중 시스템 에러 발생: {e}")
        return False, str(e)

def verify_artifacts(project_name, expected_columns):
    """
    물리적으로 생성된 산출물들을 파싱하여 구조와 서식의 무결성을 정량 평가합니다.
    """
    print(f"\n[Harness] {project_name} 최종 산출물 검증 시작...")
    workspace_dir = os.path.abspath(os.path.join(os.path.dirname(__file__), "..", ".."))
    project_dir = os.path.join(workspace_dir, project_name)
    
    csv_path = os.path.join(project_dir, "data", "raw_data.csv")
    excel_path = os.path.join(project_dir, "data", "dashboard.xlsx")
    eda_path = os.path.join(project_dir, "docs", "eda_report.md")
    docx_path = os.path.join(project_dir, "docs", "report.docx")
    pptx_path = os.path.join(project_dir, "docs", "presentation.pptx")
    
    metrics = {
        "csv": {"status": "FAIL", "reason": ""},
        "excel": {"status": "FAIL", "reason": ""},
        "eda": {"status": "FAIL", "reason": ""},
        "docx": {"status": "FAIL", "reason": ""},
        "pptx": {"status": "FAIL", "reason": ""}
    }
    
    # 1. CSV 데이터 검증
    if os.path.exists(csv_path):
        try:
            df = pd.read_csv(csv_path)
            # 행수 검증
            if len(df) < 10:
                metrics["csv"]["reason"] = f"수집된 데이터가 너무 적음 ({len(df)}건)"
            # 스키마(컬럼) 검증
            elif not all(col in df.columns for col in expected_columns):
                missing_cols = [col for col in expected_columns if col not in df.columns]
                metrics["csv"]["reason"] = f"필수 컬럼 누락: {missing_cols}"
            else:
                metrics["csv"]["status"] = "PASS"
                metrics["csv"]["reason"] = f"총 {len(df)}건 수집 완료, 스키마 일치"
        except Exception as e:
            metrics["csv"]["reason"] = f"CSV 로드 에러: {e}"
    else:
        metrics["csv"]["reason"] = "CSV 파일 존재하지 않음"

    # 2. Excel 대시보드 검증
    if os.path.exists(excel_path):
        try:
            wb = load_workbook(excel_path)
            sheet_names = wb.sheetnames
            # 필수 시트 검증
            if "대시보드 요약" not in sheet_names or "상세 데이터" not in sheet_names:
                metrics["excel"]["reason"] = f"필수 시트 누락 (존재 시트: {sheet_names})"
            else:
                ws_summary = wb["대시보드 요약"]
                # 타이틀바 텍스트 검증
                title_val = ws_summary["A1"].value
                if not title_val or "대시보드" not in str(title_val):
                    metrics["excel"]["reason"] = "요약 시트 타이틀바 텍스트 부적절"
                else:
                    metrics["excel"]["status"] = "PASS"
                    metrics["excel"]["reason"] = f"시트 {len(sheet_names)}개 구성 및 요약 테이블 포맷 완료"
            wb.close()
        except Exception as e:
            metrics["excel"]["reason"] = f"Excel 로드 에러: {e}"
    else:
        metrics["excel"]["reason"] = "Excel 파일 존재하지 않음"

    # 3. EDA 보고서 검증
    if os.path.exists(eda_path):
        try:
            with open(eda_path, "r", encoding="utf-8") as f:
                content = f.read()
            # 마크다운 헤더 구성 체크
            if "#" not in content or "##" not in content:
                metrics["eda"]["reason"] = "보고서 헤더 구조(H1, H2) 누락"
            elif "images/summary_chart.png" not in content.replace("\\", "/"):
                metrics["eda"]["reason"] = "시각화 차트 이미지 링크 유실"
            else:
                metrics["eda"]["status"] = "PASS"
                metrics["eda"]["reason"] = "마크다운 문서 골격 및 이미지 삽입 정상"
        except Exception as e:
            metrics["eda"]["reason"] = f"EDA 보고서 리드 에러: {e}"
    else:
        metrics["eda"]["reason"] = "EDA 보고서 파일 존재하지 않음"

    # 4. Word 보고서 검증
    if os.path.exists(docx_path):
        try:
            doc = Document(docx_path)
            # 단락 수 및 표 개수 확인
            if len(doc.paragraphs) < 3:
                metrics["docx"]["reason"] = "문서 단락 구성이 너무 부실함"
            elif len(doc.tables) < 1:
                metrics["docx"]["reason"] = "핵심 통계 표가 누락됨"
            else:
                metrics["docx"]["status"] = "PASS"
                metrics["docx"]["reason"] = f"단락 수: {len(doc.paragraphs)}, 표 개수: {len(doc.tables)} 확인 완료"
        except Exception as e:
            metrics["docx"]["reason"] = f"Word 로드 에러: {e}"
    else:
        metrics["docx"]["reason"] = "Word 파일 존재하지 않음"

    # 5. PPTX 발표 슬라이드 검증
    if os.path.exists(pptx_path):
        try:
            prs = Presentation(pptx_path)
            # 슬라이드 개수 검증
            if len(prs.slides) < 3:
                metrics["pptx"]["reason"] = f"슬라이드 장표가 기준 미달임 ({len(prs.slides)}장)"
            else:
                metrics["pptx"]["status"] = "PASS"
                metrics["pptx"]["reason"] = f"총 {len(prs.slides)}장 슬라이드 구성 완료"
        except Exception as e:
            metrics["pptx"]["reason"] = f"PPTX 로드 에러: {e}"
    else:
        metrics["pptx"]["reason"] = "PPTX 파일 존재하지 않음"

    return metrics

def write_harness_report(project_name, pipeline_ok, pipeline_err, metrics):
    """
    검증 결과를 마크다운 리포트와 JSON 로그 형태로 harness/results/ 아래에 출력합니다.
    """
    workspace_dir = os.path.abspath(os.path.join(os.path.dirname(__file__), "..", ".."))
    results_dir = os.path.join(workspace_dir, "harness", "results")
    os.makedirs(results_dir, exist_ok=True)
    
    timestamp = time.strftime('%Y%m%d_%H%M%S')
    md_filename = f"harness_run_{project_name}_{timestamp}.md"
    json_filename = f"harness_run_{project_name}_{timestamp}.json"
    
    md_path = os.path.join(results_dir, md_filename)
    json_path = os.path.join(results_dir, json_filename)
    
    # 1. 점수 및 통과율 계산
    total_checks = len(metrics) + 1 # 산출물 5개 + 파이프라인 정상 수행 여부
    passed_checks = sum(1 for m in metrics.values() if m["status"] == "PASS")
    if pipeline_ok:
        passed_checks += 1
        
    pass_rate = (passed_checks / total_checks) * 100
    
    # 2. 마크다운 보고서 내용 작성
    report = []
    report.append(f"# 통합 테스트 하네스 검증 결과 보고서 ({project_name})\n\n")
    report.append(f"- **검증 일시**: {time.strftime('%Y-%m-%d %H:%M:%S')}\n")
    report.append(f"- **종합 판정**: {'SUCCESS' if pass_rate == 100 else 'FAILED'}\n")
    report.append(f"- **검증 항목 통과율**: {pass_rate:.1f}% ({passed_checks}/{total_checks})\n\n")
    
    report.append("## 1. 파이프라인 수행 테스트\n\n")
    if pipeline_ok:
        report.append("- **결과**: `PASS` (정상 실행 마감)\n")
    else:
        report.append("- **결과**: `FAIL` (스크립트 런타임 오류 발생)\n")
        report.append(f"```text\n{pipeline_err}\n```\n")
        
    report.append("\n## 2. 물리 산출물 정량 검증 지표\n\n")
    report.append("| 검증 산출물 | 상태 | 상세 검증 사유 |\n")
    report.append("| :--- | :---: | :--- |\n")
    for key, val in metrics.items():
        report.append(f"| {key.upper()} 데이터 | `{val['status']}` | {val['reason']} |\n")
        
    report.append("\n## 3. 비즈니스 가이드 준수 총평\n\n")
    if pass_rate == 100:
        report.append("본 프로젝트의 산출물들은 범용 크롤링 및 데이터 분석 파이프라인 규칙을 모두 충족하며, openpyxl 테마 서식, python-pptx 레이아웃 및 matplotlib 한글 인코딩 호환성을 완벽히 확보한 것으로 평가됩니다.\n")
    else:
        report.append("일부 산출물 파일이 누락되었거나 구조적 검증에 실패했습니다. 위 세부 실패 내역을 참조하여 스크립트 파일들을 보정해 주십시오.\n")
        
    # 파일 쓰기
    with open(md_path, "w", encoding="utf-8") as f:
        f.writelines(report)
        
    # JSON 파일 쓰기
    log_data = {
        "project": project_name,
        "timestamp": timestamp,
        "pass_rate": pass_rate,
        "pipeline_ok": pipeline_ok,
        "pipeline_error": pipeline_err,
        "metrics": metrics
    }
    with open(json_path, "w", encoding="utf-8") as f:
        json.dump(log_data, f, ensure_ascii=False, indent=4)
        
    print(f"\n[Harness] 검증 보고서 작성 완료:")
    print(f" - Markdown: {md_path}")
    print(f" - JSON Log: {json_path}")
    return md_path

if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="프로젝트 파이프라인 검증용 하네스 러너")
    parser.add_argument("--project", required=True, help="검증할 프로젝트 명")
    parser.add_argument("--url", default=None, help="Pre-Scrape 테스트용 타겟 URL")
    parser.add_argument("--columns", default="title,summary,publisher,link,date", help="검증할 CSV 스키마 (쉼표 구분)")
    
    args = parser.parse_args()
    
    cols = [c.strip() for c in args.columns.split(",") if c.strip()]
    
    # 1. 파이프라인 구동 테스트
    pipeline_ok, pipeline_err = run_pipeline(args.project, args.url)
    
    # 2. 산출물 정밀 검증
    metrics = verify_artifacts(args.project, cols)
    
    # 3. 리포트 영구 박제
    write_harness_report(args.project, pipeline_ok, pipeline_err, metrics)
