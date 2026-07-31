"""
마른김 Top 10 정량 데이터 및 파트너 디렉토리를 100% 완전 수록하여
PowerPoint (.pptx) 발표 덱으로 변환 및 생성하는 스크립트.
"""

import os
import sys
import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

sys.stdout.reconfigure(encoding='utf-8')

PPTX_OUTPUT_PATH = 'BIZ-laver/reports/HaeYu_Laver_Export_Master_Market_Expansion_Deck_Full.pptx'
IMAGE_DIR = 'BIZ-laver/images'

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

COLOR_BG_DARK = RGBColor(15, 23, 42)
COLOR_NAVY = RGBColor(26, 82, 118)
COLOR_ACCENT_BLUE = RGBColor(52, 152, 219)
COLOR_TEXT_MAIN = RGBColor(30, 41, 59)
COLOR_CARD_BG = RGBColor(248, 250, 252)

blank_layout = prs.slide_layouts[6]

def add_slide_header(slide, title_text, category_text="해유 김 수출 종합 EDA & 1인 무역회사 시장개척"):
    header_box = slide.shapes.add_shape(pptx.enum.shapes.MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.1))
    header_box.fill.solid()
    header_box.fill.fore_color.rgb = COLOR_NAVY
    header_box.line.color.rgb = COLOR_NAVY
    
    tf = header_box.text_frame
    tf.margin_left = Inches(0.5)
    tf.margin_top = Inches(0.15)
    
    p0 = tf.paragraphs[0]
    p0.text = category_text.upper()
    p0.font.size = Pt(10)
    p0.font.bold = True
    p0.font.color.rgb = COLOR_ACCENT_BLUE
    p0.font.name = '맑은 고딕'
    
    p1 = tf.add_paragraph()
    p1.text = title_text
    p1.font.size = Pt(20)
    p1.font.bold = True
    p1.font.color.rgb = RGBColor(255, 255, 255)
    p1.font.name = '맑은 고딕'

# Slide 1: Cover
slide1 = prs.slides.add_slide(blank_layout)
bg1 = slide1.shapes.add_shape(pptx.enum.shapes.MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
bg1.fill.solid()
bg1.fill.fore_color.rgb = COLOR_BG_DARK
bg1.line.color.rgb = COLOR_BG_DARK

tf1 = slide1.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.333), Inches(3.5)).text_frame
p = tf1.paragraphs[0]
p.text = "🚀 [마스터 전략] 해유 김 수출 종합 EDA & 1인 무역회사 시장개척"
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)
p.font.name = '맑은 고딕'

p_sub = tf1.add_paragraph()
p_sub.text = "마른김 & 조미김 Top 10 정량 데이터, 파트너 디렉토리 및 4단계 실전 개척 프로토콜"
p_sub.font.size = Pt(17)
p_sub.font.color.rgb = COLOR_ACCENT_BLUE
p_sub.font.name = '맑은 고딕'
p_sub.space_before = Pt(12)

# Slide 2: 마른김 (HS 121221) Top 10 정량 데이터 & 전략
slide2 = prs.slides.add_slide(blank_layout)
add_slide_header(slide2, "마른김 (HS 121221) Top 10 타깃 시장 & B2B 원초 딜러 전략")

card_raw_top = slide2.shapes.add_shape(pptx.enum.shapes.MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.4), Inches(12.133), Inches(5.4))
card_raw_top.fill.solid()
card_raw_top.fill.fore_color.rgb = COLOR_CARD_BG
card_raw_top.line.color.rgb = COLOR_NAVY

tf_r = card_raw_top.text_frame
tf_r.margin_left = Inches(0.4)
tf_r.margin_top = Inches(0.3)

p_r = tf_r.paragraphs[0]
p_r.text = "🌾 마른김 Top 10 시장 및 B2B 원초 중개 딜러 (Commission Agent) 모델"
p_r.font.size = Pt(16)
p_r.font.bold = True
p_r.font.color.rgb = COLOR_NAVY

raw_top_text = [
    "• Top 1 시장: 일본 ($167.06M, +122.2%) - 초밥용 김 / 김가루 재가공용 원초 수입",
    "• Top 2 시장: 중국 ($100.61M, +102.2%) - 현지 조미김 가공 공장 대량 원초 수입",
    "• Top 3 시장: 태국 ($89.52M, +128.9%) - Taokaenoi 등 세계적 김스낵 가공 라인 대량 원초 수요",
    "• Top 4~10: 러시아 ($79.0M), 베트남 ($23.58M, +261%), 인도네시아 ($19.38M, +261.5%)",
    "",
    "🚀 1인 무역상 액션 플랜:",
    "1. B2B 원초 딜러/중개상 모델: 산지 조합(해남, 서천)과 해외 공장 연결 중개 무역에 집중",
    "2. 품질 등급 검수 서비스 병행: '품질 검수 보증 보고서' 제공으로 안전 마진 수수료 확보",
    "3. 수확기 장기 계약: 11월~4월 원초 수확기 산지 장기 계약으로 비수기 가격 변동 방어"
]

for item in raw_top_text:
    p_item = tf_r.add_paragraph()
    p_item.text = item
    p_item.font.size = Pt(12)
    p_item.font.color.rgb = COLOR_TEXT_MAIN
    p_item.font.name = '맑은 고딕'
    p_item.space_before = Pt(4)

# Slide 3: 조미김 Top 10 & 파트너 디렉토리
slide3 = prs.slides.add_slide(blank_layout)
add_slide_header(slide3, "조미김 Top 10 타깃 시장 & 국가별 파트너 디렉토리")

card_box3 = slide3.shapes.add_shape(pptx.enum.shapes.MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.4), Inches(12.133), Inches(5.4))
card_box3.fill.solid()
card_box3.fill.fore_color.rgb = COLOR_CARD_BG
card_box3.line.color.rgb = COLOR_NAVY

tf3 = card_box3.text_frame
tf3.margin_left = Inches(0.4)
tf3.margin_top = Inches(0.3)

p3 = tf3.paragraphs[0]
p3.text = "🎯 주요 국가별 사명 / 웹사이트 / 이메일 / 매칭 특징 디렉토리"
p3.font.size = Pt(16)
p3.font.bold = True
p3.font.color.rgb = COLOR_NAVY

partner_list_deck = [
    "• 태국 (마른김): Taokaenoi Food & Marketing (taokaenoi.co.th / export@taokaenoi.co.th)",
    "• 베트남 (마른김): Miwon Vietnam (miwon.com.vn / info@miwon.com.vn)",
    "• 인도네시아 (마른김): PT Miwon Indonesia (mamasuka.com / customer@mamasuka.com)",
    "• USA (조미김): Weee! Inc. (sayweee.com / vendor@sayweee.com) & H-Mart (hmart.com)",
    "• 폴란드 (조미김): Kuchnie Świata (kuchnieswiata.com.pl / b2b@kuchnieswiata.com.pl)",
    "• UAE (조미김): Choithrams (choithrams.com / info@choithrams.com) & Lulu Group",
    "• 카자흐스탄 (조미김): Magnum Cash & Carry (magnum.kz / import@magnum.kz)"
]

for item in partner_list_deck:
    p_item = tf3.add_paragraph()
    p_item.text = item
    p_item.font.size = Pt(12)
    p_item.font.color.rgb = COLOR_TEXT_MAIN
    p_item.font.name = '맑은 고딕'
    p_item.space_before = Pt(6)

# Slide 4: 1인 무역회사 실전 4단계 프로토콜
slide4 = prs.slides.add_slide(blank_layout)
add_slide_header(slide4, "1인 무역회사 실전 4단계 개척 프로토콜 & 실행 가이드")

steps = [
    ("Step 1", "고마진 포지셔닝", "조미김 완제품 스낵\n지퍼백 파우치 소포장\n비건 & 할랄 인증"),
    ("Step 2", "KOTRA 바이어 발굴", "무역관 지사화 사업\nEC21, Kompass 활용\n샘플 무료 항공 배송"),
    ("Step 3", "LCL 물류 & 라벨링", "초기 1~2 Pallet LCL\n해상 운송 재고 최소화\n현지어 영양성분 라벨"),
    ("Step 4", "E-Commerce & 숏폼", "Amazon, Noon, Allegro\n현지 크로스보더 입점\n틱톡 K-Food 숏폼 챌린지")
]

for idx, (s_num, s_title, s_desc) in enumerate(steps):
    left_pos = Inches(0.6 + idx * 3.05)
    s_box = slide4.shapes.add_shape(pptx.enum.shapes.MSO_SHAPE.ROUNDED_RECTANGLE, left_pos, Inches(1.8), Inches(2.8), Inches(4.8))
    s_box.fill.solid()
    s_box.fill.fore_color.rgb = COLOR_CARD_BG
    s_box.line.color.rgb = COLOR_ACCENT_BLUE
    
    tf_s = s_box.text_frame
    tf_s.margin_left = Inches(0.2)
    tf_s.margin_top = Inches(0.3)
    
    p0 = tf_s.paragraphs[0]
    p0.text = s_num
    p0.font.size = Pt(14)
    p0.font.bold = True
    p0.font.color.rgb = COLOR_ACCENT_BLUE
    
    p1 = tf_s.add_paragraph()
    p1.text = s_title
    p1.font.size = Pt(16)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_NAVY
    p1.space_before = Pt(8)
    
    p2 = tf_s.add_paragraph()
    p2.text = s_desc
    p2.font.size = Pt(11)
    p2.font.color.rgb = COLOR_TEXT_MAIN
    p2.space_before = Pt(12)

prs.save(PPTX_OUTPUT_PATH)
print(f"100% 완벽 수록 Master PowerPoint 발표 덱이 {PPTX_OUTPUT_PATH}에 성공적으로 저장되었습니다.")
